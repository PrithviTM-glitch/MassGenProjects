from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor

# Initialize presentation
prs = Presentation()
blank_layout = prs.slide_layouts[6] # Layout 6 is entirely blank

def add_text(slide, text, left, top, width, height, font_size=18, bold=False, center=False):
    """Helper function to create custom text boxes without bullet points."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    if center:
        p.alignment = PP_ALIGN.CENTER
    return txBox

def add_quiz_button(slide, top_position):
    """Helper function to add a visually distinct, clickable quiz button."""
    left = Inches(2.5)
    width = Inches(5)
    height = Inches(1)
    
    # Add a rounded rectangle shape
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top_position, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(0, 102, 204) # Nice blue color
    
    # Add text and hyperlink to the shape
    tf = shape.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    
    run = p.add_run()
    run.text = "🧠 QUIZ: Click Here to Answer"
    run.font.bold = True
    run.font.size = Pt(20)
    run.font.color.rgb = RGBColor(255, 255, 255) # White text
    # Hyperlink to the next slide (or you can put a URL here like "https://forms.google.com/...")
    run.hyperlink.address = "https://example.com/quiz" 

# --- Slide 1: Welcome ---
slide1 = prs.slides.add_slide(blank_layout)
add_text(slide1, "The Quantum Maze", Inches(1), Inches(2.5), Inches(8), Inches(1), font_size=44, bold=True, center=True)
add_text(slide1, "Navigating the Future of Computing.", Inches(1), Inches(3.5), Inches(8), Inches(1), font_size=24, center=True)

# --- Slide 2: The Classical Runner ---
slide2 = prs.slides.add_slide(blank_layout)
add_text(slide2, "The Classical Runner", Inches(1), Inches(0.5), Inches(8), Inches(1), font_size=32, bold=True, center=True)
add_text(slide2, "Imagine being dropped into a massive maze. A standard, classical computer acts exactly like a single, ordinary runner. It runs down one path, hits a dead end, turns around, walks back, and tries the next available route.", Inches(0.5), Inches(2), Inches(4.2), Inches(4), font_size=20)
add_text(slide2, "This trial-and-error method works perfectly fine for simple puzzles. However, if the maze is the size of a galaxy, even the fastest supercomputer in the world would take millions of years of running back and forth to finally find the exit.", Inches(5.3), Inches(2), Inches(4.2), Inches(4), font_size=20)

# --- Slide 3: Breaking the Rules ---
slide3 = prs.slides.add_slide(blank_layout)
add_text(slide3, "Throwing Away the Rulebook", Inches(1), Inches(1), Inches(8), Inches(1), font_size=36, bold=True, center=True)
add_text(slide3, "What if you did not have to choose just one path at a time? Quantum computing fundamentally changes the rules of the game. Instead of sending a single runner to test paths sequentially, it alters the very nature of how we interact with the puzzle itself.", Inches(1.5), Inches(3), Inches(7), Inches(3), font_size=24, center=True)

# --- Slide 4: Knowledge Check 1 ---
slide4 = prs.slides.add_slide(blank_layout)
add_text(slide4, "Knowledge Check", Inches(1), Inches(1.5), Inches(8), Inches(1), font_size=36, bold=True, center=True)
add_text(slide4, "Based on our rules so far, how does a classical computer try to escape the maze?", Inches(1.5), Inches(2.5), Inches(7), Inches(1.5), font_size=24, center=True)
add_quiz_button(slide4, Inches(4.5))

# --- Slide 5: The Qubit Coin ---
slide5 = prs.slides.add_slide(blank_layout)
add_text(slide5, "The Qubit Coin", Inches(1), Inches(0.5), Inches(8), Inches(1), font_size=32, bold=True, center=True)
add_text(slide5, "Normal computers use bits, which act like light switches permanently stuck on either zero or one. In our maze, a bit means you are absolutely committing to turning left, or you are absolutely committing to turning right. There is no in-between.", Inches(0.5), Inches(2), Inches(4.2), Inches(4), font_size=20)
add_text(slide5, "Quantum computers use 'qubits.' Think of a qubit like a spinning coin. While it spins in the air, it is not just heads or tails—it exists as a fluid blur of both possibilities at the exact same time.", Inches(5.3), Inches(2), Inches(4.2), Inches(4), font_size=20)

# --- Slide 6: Superposition & Entanglement ---
slide6 = prs.slides.add_slide(blank_layout)
add_text(slide6, "Superposition & Entanglement", Inches(1), Inches(0.5), Inches(8), Inches(1), font_size=32, bold=True, center=True)
add_text(slide6, "Because qubits can be multiple things at once, our maze runner achieves Superposition. The runner turns into a fluid wave, washing through the corridors and exploring every single path in the maze simultaneously.", Inches(0.5), Inches(2), Inches(4.2), Inches(4), font_size=20)
add_text(slide6, "We also have Entanglement, which acts like telepathic teammates. If two qubits become entangled, they share information instantly. If one part of the wave hits a dead end, the other side of the maze instantly knows without communicating.", Inches(5.3), Inches(2), Inches(4.2), Inches(4), font_size=20)

# --- Slide 7: Clearing the Fog ---
slide7 = prs.slides.add_slide(blank_layout)
add_text(slide7, "Quantum Interference", Inches(1), Inches(1), Inches(8), Inches(1), font_size=36, bold=True, center=True)
add_text(slide7, "Exploring every single path at once sounds incredibly chaotic. To fix this, the system uses Quantum Interference. This acts like a filter, actively canceling out all the wrong paths and dead ends while amplifying the correct route, guiding the final answer out of the maze.", Inches(1.5), Inches(3), Inches(7), Inches(3), font_size=24, center=True)

# --- Slide 8: Knowledge Check 2 ---
slide8 = prs.slides.add_slide(blank_layout)
add_text(slide8, "Knowledge Check", Inches(1), Inches(1.5), Inches(8), Inches(1), font_size=36, bold=True, center=True)
add_text(slide8, "What is the term for our runner turning into a wave and exploring all paths simultaneously?", Inches(1.5), Inches(2.5), Inches(7), Inches(1.5), font_size=24, center=True)
add_quiz_button(slide8, Inches(4.5))

# --- Slide 9: The Real-World Mazes ---
slide9 = prs.slides.add_slide(blank_layout)
add_text(slide9, "The Real-World Mazes", Inches(1), Inches(0.5), Inches(8), Inches(1), font_size=32, bold=True, center=True)
add_text(slide9, "Quantum computers will not replace your gaming laptop. They are built specifically for impossibly complex mazes, like simulating how millions of molecules interact to design new life-saving medicines in days instead of decades.", Inches(0.5), Inches(2), Inches(4.2), Inches(4), font_size=20)
add_text(slide9, "They will also revolutionize global cybersecurity. The cryptographic codes protecting our digital lives today are essentially just giant mathematical mazes. A powerful enough quantum computer could navigate those mazes and crack the codes in seconds.", Inches(5.3), Inches(2), Inches(4.2), Inches(4), font_size=20)

# --- Slide 10: The Exit is the Beginning ---
slide10 = prs.slides.add_slide(blank_layout)
add_text(slide10, "Escaping the Maze", Inches(1), Inches(1), Inches(8), Inches(1), font_size=36, bold=True, center=True)
add_text(slide10, "We are just taking our very first steps into the quantum realm. The computers being built today are early, fragile prototypes, but the runners are gearing up. The impossible mazes of tomorrow are waiting to be solved.", Inches(1.5), Inches(3), Inches(7), Inches(3), font_size=24, center=True)

# Save the presentation
prs.save('Quantum_Maze_Presentation.pptx')
print("Success! Presentation saved as 'Quantum_Maze_Presentation.pptx'")

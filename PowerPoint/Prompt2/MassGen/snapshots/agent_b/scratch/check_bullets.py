from pptx import Presentation

def check_pptx(path, label):
    print(f"\n--- Checking {label}: {path} ---")
    try:
        prs = Presentation(path)
    except Exception as e:
        print(f"Error opening {path}: {e}")
        return

    for i, slide in enumerate(prs.slides):
        print(f"Slide {i+1}:")
        for shape in slide.shapes:
            if hasattr(shape, "text_frame"):
                for p in shape.text_frame.paragraphs:
                    # Check for bullet level or characteristic bullet characters
                    has_bullet = p.level > 0 or p.text.strip().startswith(('•', '-', '*', '1.'))
                    print(f"  [Bullet: {has_bullet}] {p.text[:50]}...")

check_pptx("/Users/tmprithvi/Code/temp_workspaces/agent_b/agent1/deliverable/Quantum_Computing_Maze.pptx", "Agent 1.1")
check_pptx("/Users/tmprithvi/Code/temp_workspaces/agent_b/agent2/deliverable/Quantum_Maze_Final.pptx", "Agent 2.3")
check_pptx("/Users/tmprithvi/Code/temp_workspaces/agent_b/agent3/deliverable/Quantum_Computing_Maze_Metaphor.pptx", "Agent 3.1")

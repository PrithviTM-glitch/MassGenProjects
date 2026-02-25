from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import os

def create_quantum_pptx(output_path):
    prs = Presentation()
    
    # Define colors
    MAZE_DARK = RGBColor(20, 20, 40)
    MAZE_ACCENT = RGBColor(0, 255, 204) # Cyan/Neon
    TEXT_WHITE = RGBColor(255, 255, 255)
    QUIZ_GOLD = RGBColor(255, 215, 0)
    
    def add_maze_decoration(slide):
        # Add varied maze-like "walls" to the corners
        # Top-left corner
        wall1 = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.2), Inches(0.2), Inches(0.05), Inches(1.5))
        wall1.fill.solid()
        wall1.fill.fore_color.rgb = MAZE_ACCENT
        wall1.line.fill.background()
        
        wall2 = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.2), Inches(0.2), Inches(1.5), Inches(0.05))
        wall2.fill.solid()
        wall2.fill.fore_color.rgb = MAZE_ACCENT
        wall2.line.fill.background()
        
        # Bottom-right corner - more complex path
        wall3 = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(9.7), Inches(5.8), Inches(0.05), Inches(1.5))
        wall3.fill.solid()
        wall3.fill.fore_color.rgb = MAZE_ACCENT
        wall3.line.fill.background()
        
        wall4 = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(8.3), Inches(7.2), Inches(1.5), Inches(0.05))
        wall4.fill.solid()
        wall4.fill.fore_color.rgb = MAZE_ACCENT
        wall4.line.fill.background()
        
        # Top-right corner
        wall5 = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(8.5), Inches(0.2), Inches(1.5), Inches(0.05))
        wall5.fill.solid()
        wall5.fill.fore_color.rgb = MAZE_ACCENT
        wall5.line.fill.background()
        
        wall6 = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(9.7), Inches(0.2), Inches(0.05), Inches(0.8))
        wall6.fill.solid()
        wall6.fill.fore_color.rgb = MAZE_ACCENT
        wall6.line.fill.background()

        # Bottom-left corner
        wall7 = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.2), Inches(6.5), Inches(0.05), Inches(0.8))
        wall7.fill.solid()
        wall7.fill.fore_color.rgb = MAZE_ACCENT
        wall7.line.fill.background()
        
        wall8 = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.2), Inches(7.2), Inches(0.8), Inches(0.05))
        wall8.fill.solid()
        wall8.fill.fore_color.rgb = MAZE_ACCENT
        wall8.line.fill.background()

    def add_qubit_node(slide, left, top):
        # Add a "quantum node" (circle with glow effect)
        size = 0.4
        node = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(left), Inches(top), Inches(size), Inches(size))
        node.fill.solid()
        node.fill.fore_color.rgb = MAZE_ACCENT
        node.line.color.rgb = TEXT_WHITE
        node.line.width = Pt(1)
        
        # Connection line
        line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(left+size), Inches(top+size/2-0.02), Inches(0.8), Inches(0.04))
        line.fill.solid()
        line.fill.fore_color.rgb = MAZE_ACCENT
        line.line.fill.background()

    def set_slide_background(slide, color):
        fill = slide.background.fill
        fill.solid()
        fill.fore_color.rgb = color

    def add_text_box(slide, left, top, width, height, text, font_size=Pt(24), bold=False, color=TEXT_WHITE, align=PP_ALIGN.LEFT):
        txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
        tf = txBox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = text
        p.font.size = font_size
        p.font.bold = bold
        p.font.color.rgb = color
        p.alignment = align
        return txBox

    def add_quiz_button(slide, left, top):
        width, height = Inches(2.8), Inches(0.8)
        shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), width, height)
        shape.fill.solid()
        shape.fill.fore_color.rgb = QUIZ_GOLD
        shape.line.color.rgb = TEXT_WHITE
        shape.line.width = Pt(2)
        
        tf = shape.text_frame
        p = tf.paragraphs[0]
        p.text = "TEST YOUR INTUITION"
        p.font.size = Pt(16)
        p.font.bold = True
        p.font.color.rgb = MAZE_DARK
        p.alignment = PP_ALIGN.CENTER
        
        # Hyperlink to an external quiz
        shape.click_action.hyperlink.address = "https://example.com/quantum-maze-quiz"

    # --- SLIDE 1: TITLE ---
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide, MAZE_DARK)
    add_maze_decoration(slide)
    add_qubit_node(slide, 1, 1)
    add_qubit_node(slide, 8, 6)
    add_text_box(slide, 1, 2.5, 8, 1.2, "THE QUANTUM MAZE", Pt(60), True, MAZE_ACCENT, PP_ALIGN.CENTER)
    add_text_box(slide, 1, 4.2, 8, 1, "A Survival Guide for High School Explorers", Pt(28), False, TEXT_WHITE, PP_ALIGN.CENTER)

    # --- SLIDE 2: CLASSICAL BITS ---
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide, MAZE_DARK)
    add_maze_decoration(slide)
    add_text_box(slide, 0.5, 0.5, 9, 1, "THE CLASSICAL RUNNER", Pt(36), True, MAZE_ACCENT)
    add_text_box(slide, 0.5, 2, 4.25, 4, "Imagine a maze with only one path at a time. This is your phone and your laptop today.", Pt(22))
    add_text_box(slide, 5.25, 2, 4.25, 4, "Classical bits are like a runner who can only turn left (0) or right (1). One choice, one path.", Pt(22))

    # --- SLIDE 3: THE IMPOSSIBLE MAZE ---
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide, MAZE_DARK)
    add_maze_decoration(slide)
    add_text_box(slide, 0.5, 0.5, 9, 1, "THE MAZE GETS TOO BIG", Pt(36), True, MAZE_ACCENT)
    add_text_box(slide, 0.5, 2, 9, 2, "As mazes grow (like finding new drugs or breaking codes), the classical runner gets lost.", Pt(24))
    add_text_box(slide, 0.5, 4.5, 9, 2, "It takes trillions of years to check every dead end. We need a runner who can cheat.", Pt(24), False, RGBColor(255, 100, 100))

    # --- SLIDE 4: THE QUBIT ---
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide, MAZE_DARK)
    add_maze_decoration(slide)
    add_text_box(slide, 0.5, 0.5, 9, 1, "MEET THE GHOST RUNNER", Pt(36), True, MAZE_ACCENT)
    add_text_box(slide, 0.5, 2, 4.25, 4, "A Qubit is a 'Quantum Bit'. It doesn't pick a path until it reaches the exit.", Pt(22))
    add_text_box(slide, 5.25, 2, 4.25, 4, "It's like a cloud filling the maze, sensing all directions at once.", Pt(22))
    add_quiz_button(slide, 3.6, 6)

    # --- SLIDE 5: SUPERPOSITION ---
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide, MAZE_DARK)
    add_maze_decoration(slide)
    add_text_box(slide, 0.5, 0.5, 9, 1, "SUPERPOSITION: THE MULTI-PATH", Pt(36), True, MAZE_ACCENT)
    add_text_box(slide, 0.5, 2, 9, 3, "Superposition means the runner is everywhere at once. It's not 'left' OR 'right'—it's 'left' AND 'right'.", Pt(26))
    add_text_box(slide, 0.5, 5.5, 9, 1, "The walls don't stop the quantum runner; they just wait for the outcome.", Pt(22))

    # --- SLIDE 6: ENTANGLEMENT ---
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide, MAZE_DARK)
    add_maze_decoration(slide)
    add_text_box(slide, 0.5, 0.5, 9, 1, "ENTANGLEMENT: MAGIC LINKS", Pt(36), True, MAZE_ACCENT)
    add_text_box(slide, 0.5, 2, 4.25, 4, "If you have two quantum runners, you can 'link' them. What one does, the other knows instantly.", Pt(22))
    add_text_box(slide, 5.25, 2, 4.25, 4, "Even if they are on opposite sides of the galaxy-sized maze, they solve it together.", Pt(22))

    # --- SLIDE 7: QUANTUM GATES ---
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide, MAZE_DARK)
    add_maze_decoration(slide)
    add_text_box(slide, 0.5, 0.5, 9, 1, "GATES: REDIRECTING THE CLOUD", Pt(36), True, MAZE_ACCENT)
    add_text_box(slide, 0.5, 2, 9, 3, "Quantum Gates aren't doors. They are instructions that push the 'cloud' toward the right answer.", Pt(24))
    add_text_box(slide, 0.5, 5, 9, 2, "We manipulate probability so the ghost runner 'ends up' at the exit more often.", Pt(24))

    # --- SLIDE 8: THE EXIT ---
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide, MAZE_DARK)
    add_maze_decoration(slide)
    add_text_box(slide, 0.5, 0.5, 9, 1, "INSTANT ESCAPE", Pt(36), True, MAZE_ACCENT)
    add_text_box(slide, 0.5, 2, 4.25, 4, "While classical computers check paths one by one, quantum computers find the shortcut.", Pt(22))
    add_text_box(slide, 5.25, 2, 4.25, 4, "They don't run faster; they run smarter by seeing the whole maze from above.", Pt(22))
    add_quiz_button(slide, 3.6, 6.5)

    # --- SLIDE 9: THE PRIZE ---
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide, MAZE_DARK)
    add_maze_decoration(slide)
    add_text_box(slide, 0.5, 0.5, 9, 1, "WHY DO WE CARE?", Pt(36), True, MAZE_ACCENT)
    add_text_box(slide, 0.5, 2, 9, 2, "Decoding DNA, designing new materials, and solving climate change puzzles.", Pt(24))
    add_text_box(slide, 0.5, 4.5, 9, 2, "These are mazes that would take regular computers until the end of time to solve.", Pt(24))

    # --- SLIDE 10: YOUR TURN ---
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide, MAZE_DARK)
    add_maze_decoration(slide)
    add_text_box(slide, 1, 2.5, 8, 1.5, "EXITING THE MAZE", Pt(60), True, MAZE_ACCENT, PP_ALIGN.CENTER)
    add_text_box(slide, 1, 4.2, 8, 2, "You are the next generation of Quantum Architects. Go build the path.", Pt(28), False, TEXT_WHITE, PP_ALIGN.CENTER)

    prs.save(output_path)
    print(f"Presentation saved to {output_path}")

if __name__ == "__main__":
    output_file = "quantum_computing_maze_v2.pptx"
    create_quantum_pptx(output_file)

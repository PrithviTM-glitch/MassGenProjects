#!/usr/bin/env python3
"""Generate a 10-slide PPTX explaining quantum computing using a maze metaphor.

Constraints from user:
- 10 slides
- Continuous maze visual metaphor (maze motif on every slide)
- Superposition + entanglement introduced by slide 6
- No bullet points
- Dynamic layouts (multi text boxes)
- Hyperlinked standout 'Quiz' button on slides 4 and 8
"""

from __future__ import annotations

import argparse
import math
import os
import random
from dataclasses import dataclass
from pathlib import Path

from PIL import Image, ImageDraw
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


# ----------------------------
# Maze generation + rendering
# ----------------------------


@dataclass(frozen=True)
class Maze:
    w: int
    h: int
    right_wall: list[list[bool]]  # right_wall[y][x] is wall between (x,y) and (x+1,y)
    down_wall: list[list[bool]]  # down_wall[y][x] is wall between (x,y) and (x,y+1)


def _neighbors(x: int, y: int, w: int, h: int) -> list[tuple[int, int, str]]:
    out: list[tuple[int, int, str]] = []
    if x > 0:
        out.append((x - 1, y, "L"))
    if x < w - 1:
        out.append((x + 1, y, "R"))
    if y > 0:
        out.append((x, y - 1, "U"))
    if y < h - 1:
        out.append((x, y + 1, "D"))
    return out


def generate_maze(w: int, h: int, seed: int) -> Maze:
    rnd = random.Random(seed)
    right_wall = [[True for _ in range(w)] for __ in range(h)]
    down_wall = [[True for _ in range(w)] for __ in range(h)]
    visited = [[False for _ in range(w)] for __ in range(h)]

    stack: list[tuple[int, int]] = [(0, 0)]
    visited[0][0] = True

    while stack:
        x, y = stack[-1]
        neigh = [(nx, ny, d) for (nx, ny, d) in _neighbors(x, y, w, h) if not visited[ny][nx]]
        if not neigh:
            stack.pop()
            continue
        nx, ny, d = rnd.choice(neigh)

        if d == "R":
            right_wall[y][x] = False
        elif d == "L":
            right_wall[y][nx] = False
        elif d == "D":
            down_wall[y][x] = False
        elif d == "U":
            down_wall[ny][x] = False

        visited[ny][nx] = True
        stack.append((nx, ny))

    return Maze(w=w, h=h, right_wall=right_wall, down_wall=down_wall)


def solve_maze(maze: Maze, start: tuple[int, int] = (0, 0), goal: tuple[int, int] | None = None) -> list[tuple[int, int]]:
    if goal is None:
        goal = (maze.w - 1, maze.h - 1)

    from collections import deque

    sx, sy = start
    gx, gy = goal
    q = deque([(sx, sy)])
    prev: dict[tuple[int, int], tuple[int, int] | None] = {(sx, sy): None}

    def can_move(ax: int, ay: int, bx: int, by: int) -> bool:
        # adjacent only
        if abs(ax - bx) + abs(ay - by) != 1:
            return False
        if bx == ax + 1:  # right
            return not maze.right_wall[ay][ax]
        if bx == ax - 1:  # left
            return not maze.right_wall[ay][bx]
        if by == ay + 1:  # down
            return not maze.down_wall[ay][ax]
        if by == ay - 1:  # up
            return not maze.down_wall[by][ax]
        return False

    while q:
        x, y = q.popleft()
        if (x, y) == (gx, gy):
            break
        for nx, ny, _ in _neighbors(x, y, maze.w, maze.h):
            if (nx, ny) in prev:
                continue
            if can_move(x, y, nx, ny):
                prev[(nx, ny)] = (x, y)
                q.append((nx, ny))

    if (gx, gy) not in prev:
        return [start]

    path: list[tuple[int, int]] = []
    cur: tuple[int, int] | None = (gx, gy)
    while cur is not None:
        path.append(cur)
        cur = prev[cur]
    path.reverse()
    return path


def render_maze_background(
    maze: Maze,
    path: list[tuple[int, int]],
    step_index: int,
    total_steps: int,
    out_path: Path,
    mode: str = "normal",
) -> None:
    """Render a 1920x1080 PNG with maze walls and a progress marker.

    mode:
      - normal
      - superposition
      - entanglement
    """

    W, H = 1920, 1080
    bg = (10, 16, 32)  # deep navy
    wall = (210, 225, 255)
    glow = (79, 209, 197)  # teal
    amber = (251, 191, 36)

    img = Image.new("RGB", (W, H), bg)
    d = ImageDraw.Draw(img)

    # Maze drawing region
    margin = 110
    cell = int(min((W - 2 * margin) / maze.w, (H - 2 * margin) / maze.h))
    ox = (W - maze.w * cell) // 2
    oy = (H - maze.h * cell) // 2

    # Outer border (draw as a rectangle, then carve entrance/exit)
    left = ox
    top = oy
    right = ox + maze.w * cell
    bottom = oy + maze.h * cell

    line_w = max(3, cell // 10)

    # Outer border lines
    d.rectangle([left, top, right, bottom], outline=wall, width=line_w)

    # Carve entrance (left side of cell 0,0) and exit (right side of last cell)
    entrance_y = top + cell // 2
    exit_y = bottom - cell // 2
    d.line([left, entrance_y - cell // 4, left, entrance_y + cell // 4], fill=bg, width=line_w + 2)
    d.line([right, exit_y - cell // 4, right, exit_y + cell // 4], fill=bg, width=line_w + 2)

    # Inner walls
    for y in range(maze.h):
        for x in range(maze.w):
            x0 = ox + x * cell
            y0 = oy + y * cell
            if maze.right_wall[y][x] and x != maze.w - 1:
                d.line([x0 + cell, y0, x0 + cell, y0 + cell], fill=wall, width=line_w)
            if maze.down_wall[y][x] and y != maze.h - 1:
                d.line([x0, y0 + cell, x0 + cell, y0 + cell], fill=wall, width=line_w)

    def cell_center(cx: int, cy: int) -> tuple[int, int]:
        return (ox + cx * cell + cell // 2, oy + cy * cell + cell // 2)

    # Progress along the solved path
    if len(path) <= 1:
        t = 0
    else:
        t = (step_index) / max(1, (total_steps - 1))
    path_pos = int(round(t * (len(path) - 1)))
    path_pos = max(0, min(len(path) - 1, path_pos))

    # Draw a faint path trail up to current position
    if path_pos > 0:
        pts = [cell_center(x, y) for (x, y) in path[: path_pos + 1]]
        for wmul, col in [(5, (glow[0], glow[1], glow[2],)), (2, (amber[0], amber[1], amber[2],))]:
            # PIL ImageDraw doesn't support alpha on RGB; emulate with repeated thinner strokes
            d.line(pts, fill=col, width=max(2, line_w // wmul))

    def draw_dot(xy: tuple[int, int], r: int, color: tuple[int, int, int]) -> None:
        x, y = xy
        d.ellipse([x - r, y - r, x + r, y + r], fill=color, outline=(255, 255, 255))

    cur_cell = path[path_pos]
    cur_xy = cell_center(*cur_cell)

    if mode == "normal":
        draw_dot(cur_xy, r=max(10, line_w * 2), color=glow)

    elif mode == "superposition":
        # Two simultaneous positions near the current checkpoint
        alt_a = path[path_pos]
        alt_b = path[min(len(path) - 1, path_pos + max(2, len(path) // 10))]
        draw_dot(cell_center(*alt_a), r=max(10, line_w * 2), color=glow)
        draw_dot(cell_center(*alt_b), r=max(10, line_w * 2), color=amber)
        # Add a soft bridge between them
        d.line([cell_center(*alt_a), cell_center(*alt_b)], fill=(160, 160, 255), width=max(2, line_w // 2))

    elif mode == "entanglement":
        # Two dots in different maze regions, linked by a line
        a = path[path_pos]
        b = path[max(0, path_pos - max(3, len(path) // 8))]
        axy = cell_center(*a)
        bxy = cell_center(*b)
        d.line([axy, bxy], fill=(180, 110, 255), width=max(3, line_w // 2))
        draw_dot(axy, r=max(10, line_w * 2), color=glow)
        draw_dot(bxy, r=max(10, line_w * 2), color=(180, 110, 255))

    # Add a subtle caption stripe at the bottom
    stripe_h = 70
    d.rectangle([0, H - stripe_h, W, H], fill=(8, 12, 24))
    # A tiny progress indicator line
    prog_left = 120
    prog_right = W - 120
    py = H - stripe_h // 2
    d.line([prog_left, py, prog_right, py], fill=(50, 60, 90), width=6)
    d.line([prog_left, py, int(prog_left + (prog_right - prog_left) * t), py], fill=glow, width=6)

    out_path.parent.mkdir(parents=True, exist_ok=True)
    img.save(out_path)


# ----------------------------
# PPTX building helpers
# ----------------------------


def rgb(hex6: str) -> RGBColor:
    hex6 = hex6.lstrip("#")
    return RGBColor(int(hex6[0:2], 16), int(hex6[2:4], 16), int(hex6[4:6], 16))


def add_textbox(
    slide,
    left_in: float,
    top_in: float,
    width_in: float,
    height_in: float,
    text: str,
    font_size: int,
    color: RGBColor,
    bold: bool = False,
    align: int = PP_ALIGN.LEFT,
    font_name: str = "Calibri",
):
    box = slide.shapes.add_textbox(Inches(left_in), Inches(top_in), Inches(width_in), Inches(height_in))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.alignment = align
    p.font.name = font_name
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    return box


def add_panel(
    slide,
    left_in: float,
    top_in: float,
    width_in: float,
    height_in: float,
    fill_color: RGBColor,
    line_color: RGBColor,
    line_width_pt: float = 1.5,
    radius_shape: bool = False,
):
    shape_type = MSO_SHAPE.ROUNDED_RECTANGLE if radius_shape else MSO_SHAPE.RECTANGLE
    shp = slide.shapes.add_shape(
        shape_type,
        Inches(left_in),
        Inches(top_in),
        Inches(width_in),
        Inches(height_in),
    )
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill_color
    shp.line.color.rgb = line_color
    shp.line.width = Pt(line_width_pt)
    return shp


def add_quiz_button(slide, *, left_in: float, top_in: float, w_in: float, h_in: float, target_slide):
    # Shadow
    shadow = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(left_in + 0.06),
        Inches(top_in + 0.06),
        Inches(w_in),
        Inches(h_in),
    )
    shadow.fill.solid()
    shadow.fill.fore_color.rgb = rgb("#111827")
    shadow.line.color.rgb = rgb("#111827")

    btn = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(left_in),
        Inches(top_in),
        Inches(w_in),
        Inches(h_in),
    )
    btn.fill.solid()
    btn.fill.fore_color.rgb = rgb("#FBBF24")
    btn.line.color.rgb = rgb("#FFFFFF")
    btn.line.width = Pt(2.0)

    tf = btn.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = "Quiz"
    p.alignment = PP_ALIGN.CENTER
    p.font.name = "Calibri"
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = rgb("#111827")

    btn.click_action.target_slide = target_slide
    return btn


def build_deck(output_path: Path, assets_dir: Path, seed: int = 7) -> None:
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    maze = generate_maze(w=18, h=10, seed=seed)
    solution = solve_maze(maze)

    titles = [
        "Entering the Maze",
        "Classic Path-Finding",
        "Building the Maze Map",
        "The Exit Door (Measurement)",
        "Superposition",
        "Entanglement",
        "Quantum Shortcuts",
        "Fixing Broken Walls (Error Correction)",
        "Real Maze Materials (Hardware)",
        "Quiz: Find the Right Exit",
    ]

    # Render 10 backgrounds
    bg_paths: list[Path] = []
    for i in range(10):
        mode = "normal"
        if i == 4:
            mode = "superposition"
        if i == 5:
            mode = "entanglement"
        bg = assets_dir / f"maze_slide_{i+1:02d}.png"
        render_maze_background(maze, solution, i, 10, bg, mode=mode)
        bg_paths.append(bg)

    blank_layout = prs.slide_layouts[6]
    slides = []

    for i in range(10):
        slide = prs.slides.add_slide(blank_layout)
        slides.append(slide)

        # Background image
        slide.shapes.add_picture(str(bg_paths[i]), 0, 0, width=prs.slide_width, height=prs.slide_height)

        # Slide number
        add_textbox(
            slide,
            11.9,
            7.05,
            1.3,
            0.35,
            f"{i+1}/10",
            font_size=12,
            color=rgb("#CBD5E1"),
            bold=True,
            align=PP_ALIGN.RIGHT,
        )

    # Common colors
    ink = rgb("#0B1220")
    paper = rgb("#F8FAFC")
    paper2 = rgb("#E2E8F0")
    accent = rgb("#4FD1C5")

    # Slide 1: top-heavy title + single narrative panel
    s = slides[0]
    add_panel(s, 0.6, 0.4, 12.1, 1.2, fill_color=rgb("#0B1220"), line_color=accent, radius_shape=True)
    add_textbox(s, 0.9, 0.55, 11.5, 0.8, "Quantum Computing", 44, rgb("#F8FAFC"), bold=True)
    add_textbox(s, 0.95, 1.25, 11.4, 0.35, "A maze-story for high school minds", 18, rgb("#CBD5E1"))

    add_panel(s, 0.9, 2.0, 6.2, 4.9, fill_color=paper, line_color=rgb("#94A3B8"), radius_shape=True)
    add_textbox(
        s,
        1.2,
        2.25,
        5.6,
        4.4,
        "Imagine a huge maze with a start and an exit. A computer is a machine that follows rules to choose turns.\n\nClassical computers usually commit to one corridor at a time. Quantum computers use different rules, so their ‘maze-walk’ can look very strange, yet sometimes it finds the exit faster.",
        20,
        ink,
    )

    add_panel(s, 7.5, 2.6, 4.8, 2.4, fill_color=rgb("#0B1220"), line_color=accent, radius_shape=True)
    add_textbox(
        s,
        7.85,
        2.85,
        4.1,
        2.0,
        "Goal for today\n\nUse the maze to picture superposition, entanglement, and why measurement matters.",
        20,
        rgb("#E2E8F0"),
        bold=False,
    )

    # Slide 2: two side-by-side boxes
    s = slides[1]
    add_panel(s, 0.8, 0.6, 11.7, 0.9, fill_color=rgb("#0B1220"), line_color=accent, radius_shape=True)
    add_textbox(s, 1.1, 0.72, 11.0, 0.6, "Bits and decisions in a corridor", 34, rgb("#F8FAFC"), bold=True)

    add_panel(s, 0.9, 2.0, 5.9, 4.8, fill_color=paper, line_color=rgb("#94A3B8"), radius_shape=True)
    add_textbox(
        s,
        1.2,
        2.25,
        5.3,
        4.3,
        "A classical bit is like a sign that can point Left or Right.\n\nWhen the sign is set, you take one corridor. Your path through the maze is a chain of committed choices.",
        20,
        ink,
    )

    add_panel(s, 7.0, 2.0, 5.4, 4.8, fill_color=paper2, line_color=rgb("#64748B"), radius_shape=True)
    add_textbox(
        s,
        7.3,
        2.25,
        4.8,
        4.3,
        "A qubit is a sign that can be ‘set’ in a way that does not match everyday intuition.\n\nIt can carry information about multiple potential turns until you force it to answer.",
        20,
        ink,
    )

    # Slide 3: map-building layout (three tiles)
    s = slides[2]
    add_panel(s, 0.8, 0.6, 11.7, 0.9, fill_color=rgb("#0B1220"), line_color=accent, radius_shape=True)
    add_textbox(s, 1.1, 0.72, 11.0, 0.6, "Gates are like moving the walls", 34, rgb("#F8FAFC"), bold=True)

    add_panel(s, 0.9, 2.0, 3.9, 4.8, fill_color=paper, line_color=rgb("#94A3B8"), radius_shape=True)
    add_textbox(
        s,
        1.15,
        2.25,
        3.4,
        4.3,
        "A logic gate is a rule that updates your signs.\n\nIn maze language, it is like shifting a wall so a turn that was blocked becomes open.",
        18,
        ink,
    )

    add_panel(s, 4.95, 2.0, 3.9, 4.8, fill_color=paper2, line_color=rgb("#64748B"), radius_shape=True)
    add_textbox(
        s,
        5.2,
        2.25,
        3.4,
        4.3,
        "Classical gates keep you on a single corridor.\n\nQuantum gates rotate a qubit’s ‘direction’ so different corridors can interfere.",
        18,
        ink,
    )

    add_panel(s, 9.0, 2.0, 3.4, 4.8, fill_color=rgb("#0B1220"), line_color=accent, radius_shape=True)
    add_textbox(
        s,
        9.25,
        2.25,
        2.9,
        4.3,
        "Think of a circuit as a maze designer.\n\nEach gate is one redesign step, aiming to make the exit easier to reach.",
        18,
        rgb("#E2E8F0"),
    )

    # Slide 4: measurement + quiz button
    s = slides[3]
    add_panel(s, 0.8, 0.6, 9.8, 0.9, fill_color=rgb("#0B1220"), line_color=accent, radius_shape=True)
    add_textbox(s, 1.1, 0.72, 9.2, 0.6, "Measurement is opening a door", 34, rgb("#F8FAFC"), bold=True)

    add_panel(s, 0.9, 2.0, 6.4, 4.8, fill_color=paper, line_color=rgb("#94A3B8"), radius_shape=True)
    add_textbox(
        s,
        1.2,
        2.25,
        5.8,
        4.3,
        "In quantum computing, you do not get a neat ‘map’ of where the qubit is.\n\nMeasurement is like opening the exit door and demanding: ‘Which corridor are you really in?’\n\nAfter you open the door, the answer becomes a single concrete path.",
        20,
        ink,
    )

    add_panel(s, 7.6, 2.0, 4.8, 2.6, fill_color=rgb("#0B1220"), line_color=accent, radius_shape=True)
    add_textbox(
        s,
        7.95,
        2.25,
        4.1,
        2.1,
        "Key trade\n\nYou gain a definite result.\nYou lose the delicate quantum ‘in-between’ state.",
        20,
        rgb("#E2E8F0"),
    )

    # Slide 5: superposition
    s = slides[4]
    add_panel(s, 0.8, 0.6, 11.7, 0.9, fill_color=rgb("#0B1220"), line_color=rgb("#93C5FD"), radius_shape=True)
    add_textbox(s, 1.1, 0.72, 11.0, 0.6, "Superposition: the explorer is spread out", 30, rgb("#F8FAFC"), bold=True)

    add_panel(s, 0.9, 2.0, 5.9, 4.8, fill_color=paper, line_color=rgb("#94A3B8"), radius_shape=True)
    add_textbox(
        s,
        1.2,
        2.25,
        5.3,
        4.3,
        "Before measurement, a qubit can behave like an explorer that has not picked one corridor.\n\nIt is closer to a fog that occupies multiple routes at once, with a specific ‘amount’ in each corridor.",
        20,
        ink,
    )

    add_panel(s, 7.0, 2.0, 5.4, 4.8, fill_color=paper2, line_color=rgb("#64748B"), radius_shape=True)
    add_textbox(
        s,
        7.3,
        2.25,
        4.8,
        4.3,
        "Quantum gates can reshape the fog so wrong corridors fade and the right corridor brightens.\n\nThis is why the design of the gate sequence matters.",
        20,
        ink,
    )

    # Slide 6: entanglement
    s = slides[5]
    add_panel(s, 0.8, 0.6, 11.7, 0.9, fill_color=rgb("#0B1220"), line_color=rgb("#C084FC"), radius_shape=True)
    add_textbox(s, 1.1, 0.72, 11.0, 0.6, "Entanglement: two explorers share one rulebook", 30, rgb("#F8FAFC"), bold=True)

    add_panel(s, 0.9, 2.0, 11.5, 2.2, fill_color=paper, line_color=rgb("#94A3B8"), radius_shape=True)
    add_textbox(
        s,
        1.2,
        2.25,
        10.9,
        1.7,
        "Entanglement connects qubits so their outcomes are linked, even if they are far apart in the maze.\n\nIt is like tying two explorers with an invisible rope: tug one, and the other responds in a perfectly coordinated way.",
        20,
        ink,
    )

    add_panel(s, 0.9, 4.55, 7.0, 2.25, fill_color=paper2, line_color=rgb("#64748B"), radius_shape=True)
    add_textbox(
        s,
        1.2,
        4.8,
        6.4,
        1.8,
        "This does not let you send secret messages faster than light.\n\nIt gives you correlations that classical maze-walkers cannot copy.",
        18,
        ink,
    )

    add_panel(s, 8.2, 4.55, 4.2, 2.25, fill_color=rgb("#0B1220"), line_color=rgb("#C084FC"), radius_shape=True)
    add_textbox(
        s,
        8.45,
        4.8,
        3.7,
        1.8,
        "Why it matters\n\nEntanglement is fuel for quantum teleportation and many quantum algorithms.",
        18,
        rgb("#E2E8F0"),
    )

    # Slide 7: shortcuts
    s = slides[6]
    add_panel(s, 0.8, 0.6, 11.7, 0.9, fill_color=rgb("#0B1220"), line_color=accent, radius_shape=True)
    add_textbox(s, 1.1, 0.72, 11.0, 0.6, "Algorithms: learning the maze’s rhythm", 34, rgb("#F8FAFC"), bold=True)

    add_panel(s, 0.9, 2.0, 7.2, 4.8, fill_color=paper, line_color=rgb("#94A3B8"), radius_shape=True)
    add_textbox(
        s,
        1.2,
        2.25,
        6.6,
        4.3,
        "A quantum algorithm is not magic. It is a carefully chosen set of gates that makes the right exit more likely.\n\nIn maze terms, you are not checking every corridor. You are setting up wave-like patterns so bad routes cancel and good routes reinforce.",
        20,
        ink,
    )

    add_panel(s, 8.4, 2.0, 4.0, 4.8, fill_color=rgb("#0B1220"), line_color=accent, radius_shape=True)
    add_textbox(
        s,
        8.65,
        2.25,
        3.5,
        4.3,
        "Example intuition\n\nGrover-style search is like tapping the maze to amplify the corridor that hides the prize.",
        18,
        rgb("#E2E8F0"),
    )

    # Slide 8: error correction + quiz button
    s = slides[7]
    add_panel(s, 0.8, 0.6, 9.8, 0.9, fill_color=rgb("#0B1220"), line_color=accent, radius_shape=True)
    add_textbox(s, 1.1, 0.72, 9.2, 0.6, "Noise: when the maze shakes", 34, rgb("#F8FAFC"), bold=True)

    add_panel(s, 0.9, 2.0, 6.6, 4.8, fill_color=paper, line_color=rgb("#94A3B8"), radius_shape=True)
    add_textbox(
        s,
        1.2,
        2.25,
        6.0,
        4.3,
        "Real qubits are fragile. Heat, vibration, and stray electromagnetic fields can nudge the explorer into the wrong corridor.\n\nQuantum error correction is like using extra explorers to detect wall cracks without opening the exit door too early.",
        20,
        ink,
    )

    add_panel(s, 7.8, 2.0, 4.4, 2.8, fill_color=paper2, line_color=rgb("#64748B"), radius_shape=True)
    add_textbox(
        s,
        8.1,
        2.25,
        3.8,
        2.3,
        "Important idea\n\nYou protect information by spreading it across several physical qubits.",
        18,
        ink,
    )

    # Slide 9: hardware columns
    s = slides[8]
    add_panel(s, 0.8, 0.6, 11.7, 0.9, fill_color=rgb("#0B1220"), line_color=accent, radius_shape=True)
    add_textbox(s, 1.1, 0.72, 11.0, 0.6, "What are qubits made of?", 34, rgb("#F8FAFC"), bold=True)

    cols = [
        (0.9, "Photons", "Light particles can carry qubits through optical mazes. They travel fast and resist some noise."),
        (4.65, "Trapped ions", "Charged atoms in electromagnetic ‘corridors’ can store qubits with high precision, but operations can be slower."),
        (8.4, "Superconducting circuits", "Tiny circuits chilled near absolute zero can act like artificial atoms. They are fast and widely engineered."),
    ]
    for x0, head, body in cols:
        add_panel(s, x0, 2.0, 3.75, 4.8, fill_color=paper if x0 < 5 else paper2, line_color=rgb("#94A3B8"), radius_shape=True)
        add_textbox(s, x0 + 0.25, 2.2, 3.25, 0.5, head, 22, ink, bold=True, align=PP_ALIGN.CENTER)
        add_textbox(s, x0 + 0.3, 2.8, 3.15, 3.8, body, 18, ink)

    # Slide 10: quiz
    s = slides[9]
    add_panel(s, 0.8, 0.6, 11.7, 0.9, fill_color=rgb("#0B1220"), line_color=rgb("#FBBF24"), radius_shape=True)
    add_textbox(s, 1.1, 0.72, 11.0, 0.6, "Quiz: choose the right corridor", 34, rgb("#F8FAFC"), bold=True)

    q_panels = [
        (0.9, 2.0, 5.75, 2.0, "If a bit is a committed turn, what is a qubit like in the maze story?"),
        (6.65, 2.0, 5.75, 2.0, "What does ‘measurement’ do to the maze explorer?"),
        (0.9, 4.3, 5.75, 2.4, "In one sentence, describe entanglement using the rope idea."),
        (6.65, 4.3, 5.75, 2.4, "Why can noise ruin a quantum route more easily than a classical route?"),
    ]
    for l, t, w, h, txt in q_panels:
        add_panel(s, l, t, w, h, fill_color=paper, line_color=rgb("#94A3B8"), radius_shape=True)
        add_textbox(s, l + 0.3, t + 0.25, w - 0.6, h - 0.5, txt, 18, ink)

    # Back button on quiz slide
    shadow = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.96), Inches(6.9), Inches(3.2), Inches(0.55))
    shadow.fill.solid()
    shadow.fill.fore_color.rgb = rgb("#111827")
    shadow.line.color.rgb = rgb("#111827")

    back = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.9), Inches(6.84), Inches(3.2), Inches(0.55))
    back.fill.solid()
    back.fill.fore_color.rgb = rgb("#4FD1C5")
    back.line.color.rgb = rgb("#FFFFFF")
    back.line.width = Pt(2)
    tf = back.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = "Back to Slide 1"
    p.alignment = PP_ALIGN.CENTER
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = rgb("#0B1220")
    back.click_action.target_slide = slides[0]

    # Add quiz buttons on slides 4 and 8 (1-indexed: slides[3] and slides[7]) linking to slide 10
    add_quiz_button(slides[3], left_in=10.9, top_in=0.55, w_in=1.7, h_in=0.7, target_slide=slides[9])
    add_quiz_button(slides[7], left_in=10.9, top_in=0.55, w_in=1.7, h_in=0.7, target_slide=slides[9])

    # Save
    output_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(output_path))


def main() -> None:
    ap = argparse.ArgumentParser()
    ap.add_argument("--out", default="deliverable/Quantum_Computing_Maze_Metaphor.pptx")
    ap.add_argument("--assets", default="deliverable/assets")
    ap.add_argument("--seed", type=int, default=7)
    args = ap.parse_args()

    build_deck(output_path=Path(args.out), assets_dir=Path(args.assets), seed=args.seed)


if __name__ == "__main__":
    main()

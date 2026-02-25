#!/usr/bin/env python3
"""Verify the generated Quantum Computing Maze PPTX meets key constraints."""

from __future__ import annotations

import argparse
from pathlib import Path

from pptx import Presentation
from pptx.oxml.ns import qn


def paragraph_has_bullets(paragraph) -> bool:
    """Detect bulleting in a paragraph by inspecting its pPr."""
    pPr = paragraph._p.get_or_add_pPr()
    # Any of these indicates bullet styling.
    for tag in ("a:buChar", "a:buAutoNum", "a:buBlip"):
        if pPr.find(qn(tag)) is not None:
            return True
    return False


def iter_all_paragraphs(prs: Presentation):
    for slide in prs.slides:
        for shape in slide.shapes:
            if not hasattr(shape, "text_frame"):
                continue
            tf = shape.text_frame
            for p in tf.paragraphs:
                yield slide, shape, p


def find_shapes_with_text(slide, needle: str):
    out = []
    for shape in slide.shapes:
        if not hasattr(shape, "text"):
            continue
        if (shape.text or "").strip() == needle:
            out.append(shape)
    return out


def main() -> None:
    ap = argparse.ArgumentParser()
    ap.add_argument("pptx_path", nargs="?", default="deliverable/Quantum_Computing_Maze_Metaphor.pptx")
    args = ap.parse_args()

    pptx_path = Path(args.pptx_path)
    if not pptx_path.exists():
        raise SystemExit(f"Missing PPTX: {pptx_path}")

    prs = Presentation(str(pptx_path))

    # 1) slide count
    assert len(prs.slides) == 10, f"Expected 10 slides, found {len(prs.slides)}"

    # 2) no bullets
    bullet_hits = []
    for slide, shape, p in iter_all_paragraphs(prs):
        if paragraph_has_bullets(p):
            bullet_hits.append((prs.slides.index(slide) + 1, getattr(shape, "name", "<shape>"), p.text))
    assert not bullet_hits, f"Found bullet formatting in paragraphs: {bullet_hits[:5]}"

    # 3) quiz buttons exist and hyperlink internally to slide 10
    slide4 = prs.slides[3]
    slide8 = prs.slides[7]
    slide10 = prs.slides[9]

    quiz4 = find_shapes_with_text(slide4, "Quiz")
    quiz8 = find_shapes_with_text(slide8, "Quiz")

    assert quiz4, "Slide 4 missing a shape with text 'Quiz'"
    assert quiz8, "Slide 8 missing a shape with text 'Quiz'"

    def assert_targets(slide_num: int, shapes):
        ok = False
        for shp in shapes:
            ca = shp.click_action
            if ca is not None and ca.target_slide is not None and ca.target_slide == slide10:
                ok = True
        assert ok, f"Slide {slide_num} has 'Quiz' text but no internal link to slide 10"

    assert_targets(4, quiz4)
    assert_targets(8, quiz8)

    print("PASS: slide count=10; no bullet formatting; Quiz buttons on slides 4 and 8 link to slide 10")


if __name__ == "__main__":
    main()

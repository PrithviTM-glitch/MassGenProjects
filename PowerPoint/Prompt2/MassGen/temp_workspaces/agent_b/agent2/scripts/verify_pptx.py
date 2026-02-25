from pptx import Presentation

def verify_pptx(path):
    prs = Presentation(path)
    print(f"Slide count: {len(prs.slides)}")
    
    for i, slide in enumerate(prs.slides):
        print(f"Slide {i+1}:")
        # Check for bullets
        for shape in slide.shapes:
            if hasattr(shape, "text_frame"):
                for paragraph in shape.text_frame.paragraphs:
                    if paragraph.level > 0 or (hasattr(paragraph, 'bullet') and paragraph.bullet):
                        print(f"  WARNING: Bullet found in slide {i+1}")
        
        # Check for Quiz button on slide 4 and 8 (0-indexed 3 and 7)
        if i in [3, 7]:
            found_quiz = False
            for shape in slide.shapes:
                if shape.has_text_frame and "QUIZ" in shape.text_frame.text:
                    found_quiz = True
                    break
            if found_quiz:
                print(f"  PASS: Quiz button found on slide {i+1}")
            else:
                print(f"  FAIL: Quiz button NOT found on slide {i+1}")

if __name__ == "__main__":
    verify_pptx("Quantum_Maze_Final.pptx")

import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

def create_pptx():
    prs = Presentation()
    
    # Define slide size for 16:9
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    def add_background(slide, image_path):
        if os.path.exists(image_path):
            slide.shapes.add_picture(image_path, 0, 0, width=prs.slide_width, height=prs.slide_height)
        else:
            # Fallback color background if image missing
            background = slide.background
            fill = background.fill
            fill.solid()
            fill.fore_color.rgb = RGBColor(13, 27, 42) # Dark Navy

    def add_text_box(slide, text, left, top, width, height, font_size=20, color=RGBColor(255, 255, 255), bold=False, align=PP_ALIGN.LEFT):
        txBox = slide.shapes.add_textbox(left, top, width, height)
        tf = txBox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = text
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.bold = bold
        p.alignment = align
        return txBox

    def add_quiz_button(slide, left, top, target_slide):
        # Create a gold rounded rectangle
        shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, Inches(2.2), Inches(0.6))
        shape.fill.solid()
        shape.fill.fore_color.rgb = RGBColor(255, 215, 0) # Gold
        shape.line.color.rgb = RGBColor(255, 255, 255)
        
        # Add text
        tf = shape.text_frame
        p = tf.paragraphs[0]
        p.text = "QUIZ CORNER"
        p.font.size = Pt(14)
        p.font.color.rgb = RGBColor(0, 0, 0)
        p.font.bold = True
        p.alignment = PP_ALIGN.CENTER
        
        # Link to target slide
        click_action = shape.click_action
        click_action.target_slide = target_slide
        
        return shape

    content = [
        {
            "title": "THE QUANTUM MAZE",
            "layout": "top-heavy",
            "image": "assets/image_0.png",
            "left_text": "A Journey into Future Computing",
            "right_text": "Welcome to the entrance of the most complex maze in science. Today, we leave behind the simple paths of classical computers and step into the world of Quantum."
        },
        {
            "title": "ONE WAY IN: CLASSICAL BITS",
            "layout": "side-by-side",
            "image": "assets/image_1.png",
            "left_text": "Classical computers use bits—0 or 1. Think of it as a straight corridor. You can only be in one spot at a time. The path is certain, but limited.",
            "right_text": "Every choice is a simple yes or no. In our maze, this means you can only explore one hallway at a time, making complex puzzles take forever to solve."
        },
        {
            "title": "BREAKING THE WALLS",
            "layout": "top-heavy",
            "image": "assets/image_2.png",
            "left_text": "Traditional computers are hitting a physical wall. As components get smaller, they start acting weird. We need a new way to navigate the information maze.",
            "right_text": "Quantum computing doesn't just walk faster; it changes the rules of the maze entirely, allowing us to find exits that were previously hidden behind solid walls."
        },
        {
            "title": "THE FORK: MEET THE QUBIT",
            "layout": "side-by-side",
            "image": "assets/image_3.png",
            "left_text": "A Quantum Bit, or Qubit, is special. Instead of just 0 or 1, it can be both at the same time. In our maze, imagine arriving at a fork and walking down both paths simultaneously.",
            "right_text": "This isn't magic; it's physics! By being in multiple places at once, the qubit starts to map the maze much faster than any classical bit ever could.",
            "quiz": True
        },
        {
            "title": "THE MISTY HALLWAY: SUPERPOSITION",
            "layout": "top-heavy",
            "image": "assets/image_4.png",
            "left_text": "Superposition is the state of being in multiple paths at once. It's like a misty hallway where you exist as a cloud of possibilities until someone looks at you.",
            "right_text": "Measurement 'collapses' the mist. Once we check the qubit, it picks one path. The goal is to stay in the mist as long as possible to solve the puzzle."
        },
        {
            "title": "THE LINKED ROOMS: ENTANGLEMENT",
            "layout": "side-by-side",
            "image": "assets/image_5.png",
            "left_text": "Entanglement links two qubits together, no matter how far apart they are in the maze. If you change the direction of one, the other changes instantly.",
            "right_text": "This is what Einstein called 'spooky action at a distance.' It allows quantum computers to coordinate different parts of the maze with perfect synchronicity."
        },
        {
            "title": "ECHOES: QUANTUM INTERFERENCE",
            "layout": "top-heavy",
            "image": "assets/image_6.png",
            "left_text": "How do we find the right exit? We use interference. We make the wrong paths cancel each other out (like silence) and the right paths get louder (like a glow).",
            "right_text": "By carefully timing the 'echoes' of our qubits, we ensure the maze leads us exactly where we need to go, filtering out millions of dead ends instantly."
        },
        {
            "title": "THE SLIDING DOORS: GATES",
            "layout": "side-by-side",
            "image": "assets/image_7.png",
            "left_text": "Quantum Gates are the sliding doors of our maze. They don't just open or close; they rotate and shift our path, allowing us to manipulate superposition.",
            "right_text": "By combining these gates, we create Quantum Algorithms—mathematical maps that guide us through the most complex mazes in the universe.",
            "quiz": True
        },
        {
            "title": "CRUMBLING WALLS: DECOHERENCE",
            "layout": "top-heavy",
            "left_text": "The hardest part? Keeping the maze stable. Heat, light, or even a tiny vibration can cause 'decoherence,' where the quantum paths crumble back into classical ones.",
            "right_text": "Scientists work at temperatures colder than outer space to keep the maze walls strong. Protecting these fragile paths is the greatest engineering challenge of our time.",
            "image": "assets/image_8.png"
        },
        {
            "title": "THE EXIT: THE QUANTUM FUTURE",
            "layout": "side-by-side",
            "image": "assets/image_9.png",
            "left_text": "We are just at the exit of the first quantum maze. In the future, these computers will design new medicines, crack impossible codes, and solve climate change.",
            "right_text": "The maze is huge, and we've only just begun to explore. Now that you know the rules, are you ready to become a quantum navigator?"
        }
    ]

    # First pass to create all slides
    slides = []
    for i, slide_data in enumerate(content):
        slide_layout = prs.slide_layouts[6] # Blank
        slide = prs.slides.add_slide(slide_layout)
        slides.append(slide)

    # Second pass to add content
    for i, (slide, slide_data) in enumerate(zip(slides, content)):
        # Add Background
        image_path = slide_data["image"]
        add_background(slide, image_path)
        
        # Semi-transparent overlay for text readability
        overlay = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
        overlay.fill.solid()
        overlay.fill.fore_color.rgb = RGBColor(0, 0, 0)
        overlay.fill.transparency = 0.5
        overlay.line.fill.background()

        # Add Title (Top-heavy)
        add_text_box(slide, slide_data["title"], Inches(0.5), Inches(0.5), prs.slide_width - Inches(1), Inches(1), font_size=44, bold=True, color=RGBColor(0, 255, 255))

        if slide_data["layout"] == "side-by-side":
            # Left box
            add_text_box(slide, slide_data["left_text"], Inches(0.5), Inches(2), prs.slide_width/2 - Inches(0.75), Inches(4), font_size=24)
            # Right box
            add_text_box(slide, slide_data["right_text"], prs.slide_width/2 + Inches(0.25), Inches(2), prs.slide_width/2 - Inches(0.75), Inches(4), font_size=24)
        else:
            # Top-heavy layout
            add_text_box(slide, slide_data["left_text"], Inches(0.5), Inches(2), prs.slide_width - Inches(1), Inches(2), font_size=24)
            add_text_box(slide, slide_data["right_text"], Inches(0.5), Inches(4), prs.slide_width - Inches(1), Inches(2), font_size=24)

        if slide_data.get("quiz"):
            # Target is the last slide
            add_quiz_button(slide, prs.slide_width - Inches(2.7), prs.slide_height - Inches(1.2), slides[-1])

    # Save
    prs.save("Quantum_Maze_Final.pptx")
    print("PPTX saved successfully.")

if __name__ == "__main__":
    create_pptx()

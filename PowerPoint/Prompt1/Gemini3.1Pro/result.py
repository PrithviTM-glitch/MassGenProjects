import sys
from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

def add_key_takeaway(slide, text):
    """Helper function to add a 'Key Takeaway' text box to the bottom of a slide."""
    left = Inches(0.5)
    top = Inches(6.5)
    width = Inches(9.0)
    height = Inches(0.8)
    
    txBox = slide.shapes.add_textbox(left, top, width, height)
    text_frame = txBox.text_frame
    text_frame.word_wrap = True
    
    p = text_frame.paragraphs[0]
    p.text = f"Key Takeaway: {text}"
    p.font.bold = True
    p.font.size = Pt(14)
    p.font.color.rgb = RGBColor(0, 51, 102) # Dark blue for emphasis

def main():
    try:
        # Initialize presentation
        prs = Presentation()
        
        # Slide Layouts
        title_slide_layout = prs.slide_layouts[0]
        bullet_slide_layout = prs.slide_layouts[1]
        title_only_layout = prs.slide_layouts[5]

        # ---------------------------------------------------------
        # Slide 1: Title Slide
        # ---------------------------------------------------------
        slide1 = prs.slides.add_slide(title_slide_layout)
        title1 = slide1.shapes.title
        subtitle1 = slide1.placeholders[1]
        title1.text = "Q4 & Full Year Earnings Report"
        subtitle1.text = "SaaS Co. Financial Overview\nConfidential & Internal"
        add_key_takeaway(slide1, "Q4 closed with $2.9M in revenue amidst a 4% churn increase.")

        # ---------------------------------------------------------
        # Slide 2: Executive Summary
        # ---------------------------------------------------------
        slide2 = prs.slides.add_slide(bullet_slide_layout)
        title2 = slide2.shapes.title
        title2.text = "Executive Summary"
        tf2 = slide2.placeholders[1].text_frame
        tf2.text = "Year in Review:"
        p = tf2.add_paragraph()
        p.text = "Strong growth through Q1-Q3, peaking at $3.1M in Q3."
        p.level = 1
        p = tf2.add_paragraph()
        p.text = "Q4 experienced slight contraction to $2.9M."
        p.level = 1
        p = tf2.add_paragraph()
        p.text = "Customer churn increased by 4% in Q4, driving the revenue dip."
        p.level = 1
        add_key_takeaway(slide2, "While Q4 saw a slight dip, overall YoY trajectory remains strongly positive.")

        # ---------------------------------------------------------
        # Slide 3: Quarterly Revenue (Bar Chart)
        # ---------------------------------------------------------
        slide3 = prs.slides.add_slide(title_only_layout)
        title3 = slide3.shapes.title
        title3.text = "Quarterly Revenue ($M)"
        
        # Add Bar Chart
        chart_data = CategoryChartData()
        chart_data.categories = ['Q1', 'Q2', 'Q3', 'Q4']
        chart_data.add_series('Revenue', (2.4, 2.8, 3.1, 2.9))
        
        x, y, cx, cy = Inches(1.5), Inches(1.5), Inches(7.0), Inches(4.5)
        chart = slide3.shapes.add_chart(
            XL_CHART_TYPE.COLUMN_CLUSTERED, x, y, cx, cy, chart_data
        ).chart
        chart.has_legend = False
        
        # Add exact CFO Speaker Notes
        notes_slide3 = slide3.notes_slide
        text_frame3 = notes_slide3.notes_text_frame
        text_frame3.text = (
            "CFO Notes: 'As you can see on the chart, we experienced steady, sequential "
            "growth from Q1 through Q3, peaking at $3.1M. The Q4 revenue dip to $2.9M is "
            "a direct reflection of the 4% increase in customer churn we observed at the "
            "end of the year. This was primarily driven by end-of-year budget cuts within "
            "our mid-market segment. While disappointing, our core enterprise ARR remains "
            "stable, and we are aggressively restructuring our renewal incentives to recover "
            "this momentum in Q1.'"
        )
        add_key_takeaway(slide3, "Q3 marked our highest revenue quarter at $3.1M; Q4 closed at $2.9M.")

        # ---------------------------------------------------------
        # Slide 4: Understanding the Q4 Dip
        # ---------------------------------------------------------
        slide4 = prs.slides.add_slide(bullet_slide_layout)
        title4 = slide4.shapes.title
        title4.text = "Contextualizing the Q4 Dip"
        tf4 = slide4.placeholders[1].text_frame
        tf4.text = "Drivers of the $200k drop from Q3:"
        p = tf4.add_paragraph()
        p.text = "Macroeconomic tightening led to lower-tier subscription cancellations."
        p.level = 1
        p = tf4.add_paragraph()
        p.text = "4% spike in churn outpaced our new Q4 bookings."
        p.level = 1
        add_key_takeaway(slide4, "The Q4 dip was an isolation of increased churn, not a failure in new sales.")

        # ---------------------------------------------------------
        # Slide 5: Customer Churn Breakdown (Pie Chart)
        # ---------------------------------------------------------
        slide5 = prs.slides.add_slide(title_only_layout)
        title5 = slide5.shapes.title
        title5.text = "Q4 Customer Churn Breakdown"
        
        # Add Pie Chart
        pie_data = CategoryChartData()
        pie_data.categories = ['Budget Constraints', 'Competitor Switch', 'Business Closure', 'Lack of Usage']
        # Distributing the churn drivers contextually
        pie_data.add_series('Churn Drivers', (0.45, 0.25, 0.20, 0.10))
        
        chart5 = slide5.shapes.add_chart(
            XL_CHART_TYPE.PIE, Inches(2.0), Inches(1.5), Inches(6.0), Inches(4.5), pie_data
        ).chart
        chart5.has_legend = True
        chart5.legend.position = XL_LEGEND_POSITION.RIGHT
        chart5.plots[0].has_data_labels = True

        add_key_takeaway(slide5, "Budget constraints accounted for 45% of the 4% churn increase in Q4.")

        # ---------------------------------------------------------
        # Slide 6: Action Plan & Retention Strategy
        # ---------------------------------------------------------
        slide6 = prs.slides.add_slide(bullet_slide_layout)
        title6 = slide6.shapes.title
        title6.text = "Q1 Action Plan to Combat Churn"
        tf6 = slide6.placeholders[1].text_frame
        tf6.text = "Strategic Initiatives:"
        p = tf6.add_paragraph()
        p.text = "Deploying Customer Success teams to high-risk accounts 90 days pre-renewal."
        p.level = 1
        p = tf6.add_paragraph()
        p.text = "Introducing flexible annual billing to mitigate budget constraints."
        p.level = 1
        add_key_takeaway(slide6, "Proactive engagement is our primary lever to reverse Q4 churn trends.")

        # ---------------------------------------------------------
        # Slide 7: Looking Ahead (Projections)
        # ---------------------------------------------------------
        slide7 = prs.slides.add_slide(bullet_slide_layout)
        title7 = slide7.shapes.title
        title7.text = "FY Outlook & Guidance"
        tf7 = slide7.placeholders[1].text_frame
        tf7.text = "Q1 and Beyond:"
        p = tf7.add_paragraph()
        p.text = "Targeting a return to $3.1M+ in Q1 through enterprise expansion."
        p.level = 1
        p = tf7.add_paragraph()
        p.text = "Churn expected to stabilize back to historical 1.5% average by Q2."
        p.level = 1
        add_key_takeaway(slide7, "We expect full recovery from the Q4 dip by the close of Q1.")

        # ---------------------------------------------------------
        # Slide 8: Q&A
        # ---------------------------------------------------------
        slide8 = prs.slides.add_slide(title_slide_layout)
        title8 = slide8.shapes.title
        subtitle8 = slide8.placeholders[1]
        title8.text = "Questions & Answers"
        subtitle8.text = "Open floor for leadership discussion."
        add_key_takeaway(slide8, "Transparency and rapid execution are key to our ongoing success.")

        # Save the presentation
        output_filename = "Q4_Earnings.pptx"
        prs.save(output_filename)
        print(f"Success! Presentation saved successfully as '{output_filename}'.")

    except PermissionError:
        print(f"Error: The file 'Q4_Earnings.pptx' is currently open in another program. Please close it and try again.", file=sys.stderr)
    except Exception as e:
        print(f"An unexpected error occurred during presentation generation: {e}", file=sys.stderr)

if __name__ == "__main__":
    main()

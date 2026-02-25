import os
import sys
from datetime import datetime

try:
    from pptx import Presentation
    from pptx.chart.data import CategoryChartData
    from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION
    from pptx.util import Inches, Pt
    from pptx.enum.text import PP_ALIGN
    from pptx.dml.color import RGBColor
except ImportError:
    print("Error: python-pptx not found. Please install it using 'pip install python-pptx'.")
    sys.exit(1)

# Constants for Branding
PRIMARY_COLOR = RGBColor(0, 51, 102)  # Deep Indigo/Navy
SECONDARY_COLOR = RGBColor(200, 200, 200) # Light Grey
TEXT_COLOR = RGBColor(33, 33, 33)

def add_styled_takeaway(slide, text):
    """Adds a professionally styled 'Key Takeaway' text box to the slide."""
    left = Inches(0.5)
    top = Inches(6.8)
    width = Inches(9)
    height = Inches(0.5)
    
    # Add a light background box for the takeaway
    shape = slide.shapes.add_textbox(left, top, width, height)
    tf = shape.text_frame
    tf.word_wrap = True
    p = tf.add_paragraph()
    p.text = f"Key Takeaway: {text}"
    p.font.bold = True
    p.font.size = Pt(14)
    p.font.color.rgb = PRIMARY_COLOR
    p.alignment = PP_ALIGN.LEFT

def create_earnings_deck():
    """Generates an 8-slide Quarterly Earnings deck with native charts and professional notes."""
    try:
        prs = Presentation()
        # Set 16:9 aspect ratio
        prs.slide_width = Inches(13.33)
        prs.slide_height = Inches(7.5)

        # Data Points
        QUARTERS = ['Q1', 'Q2', 'Q3', 'Q4']
        REVENUE = [2.4, 2.8, 3.1, 2.9] # in $M
        CHURN_INCREASE = "4%"
        
        # --- Slide 1: Title Slide ---
        slide = prs.slides.add_slide(prs.slide_layouts[0])
        title = slide.shapes.title
        subtitle = slide.placeholders[1]
        
        title.text = "Q4 Quarterly Earnings Review"
        subtitle.text = f"SaaS Company Performance Analysis\nDate: {datetime.now().strftime('%B %Y')}\nConfidential & Proprietary"
        
        # Style Title
        for paragraph in title.text_frame.paragraphs:
            paragraph.font.color.rgb = PRIMARY_COLOR
            
        add_styled_takeaway(slide, "Fiscal year concluded with strong top-line momentum despite Q4 seasonal variance.")
        
        # Slide 1 Notes
        slide.notes_slide.notes_text_frame.text = (
            "CFO Notes: Good afternoon everyone. Today we are reviewing our Q4 and full-year performance. "
            "While we saw some headwinds in December, the overall trajectory of the business remains healthy."
        )

        # --- Slide 2: Executive Summary ---
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = "Executive Summary: FY Performance"
        body = slide.placeholders[1].text_frame
        body.text = "Summary of Key Financial Outcomes"
        p = body.add_paragraph()
        p.text = "• Annual Revenue: $11.2M (Total across all quarters)"
        p = body.add_paragraph()
        p.text = f"• Q4 Performance: $2.9M Revenue with {CHURN_INCREASE} churn uptick"
        p = body.add_paragraph()
        p.text = "• Strategic Focus: Customer retention and enterprise segment expansion"
        
        add_styled_takeaway(slide, "Company reached double-digit annual revenue for the first time.")
        
        slide.notes_slide.notes_text_frame.text = (
            "CFO Notes: The executive summary highlights our transition into a $10M+ ARR business. "
            "Our primary narrative for this quarter is 'Scaling through Complexity'."
        )

        # --- Slide 3: Revenue Performance (BAR CHART) ---
        slide = prs.slides.add_slide(prs.slide_layouts[5]) # Title Only
        slide.shapes.title.text = "Quarterly Revenue Trajectory ($M)"
        
        chart_data = CategoryChartData()
        chart_data.categories = QUARTERS
        chart_data.add_series('Revenue', REVENUE)

        x, y, cx, cy = Inches(1.5), Inches(1.5), Inches(10), Inches(4.5)
        graphic_frame = slide.shapes.add_chart(
            XL_CHART_TYPE.COLUMN_CLUSTERED, x, y, cx, cy, chart_data
        )
        chart = graphic_frame.chart
        chart.has_legend = False
        
        # Chart Formatting
        plot = chart.plots[0]
        plot.has_data_labels = True
        data_labels = plot.data_labels
        data_labels.font.size = Pt(12)
        data_labels.font.bold = True
        
        # CFO Notes for Revenue Dip
        slide.notes_slide.notes_text_frame.text = (
            "CFO Notes: Moving to Slide 3, you'll see our quarterly revenue comparison. Q4 revenue came in at $2.9M, "
            "a slight sequential dip from the $3.1M reported in Q3. This 6.4% decrease is primarily attributed to "
            "extended procurement cycles in our enterprise segment during the holiday season and a deliberate shift "
            "in professional services timing. It is important to note that YoY, this is a 20% increase over the previous Q4."
        )
        
        add_styled_takeaway(slide, "Sequential dip in Q4 is a timing-related variance; YoY growth remains the lead indicator.")

        # --- Slide 4: Growth Drivers ---
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = "Operational Growth Drivers"
        body = slide.placeholders[1].text_frame
        body.text = "Contributing Factors to Annual Success"
        body.add_paragraph().text = "• 25% increase in Average Contract Value (ACV)"
        body.add_paragraph().text = "• Strong adoption of the new 'Premium+' tier"
        body.add_paragraph().text = "• 90% Net Revenue Retention (NRR) excluding Q4 outliers"
        
        add_styled_takeaway(slide, "ACV growth is outpacing customer count growth, indicating move up-market.")
        
        slide.notes_slide.notes_text_frame.text = "CFO Notes: Our unit economics remain strong. The Premium+ tier is seeing 40% attach rates on new deals."

        # --- Slide 5: Churn Analysis (PIE CHART) ---
        slide = prs.slides.add_slide(prs.slide_layouts[5])
        slide.shapes.title.text = f"Q4 Churn Analysis: {CHURN_INCREASE} Increase Context"
        
        chart_data = CategoryChartData()
        chart_data.categories = ['Competitive Pressure', 'Budget Cuts', 'Product Fit', 'M&A/Dissolution']
        chart_data.add_series('Churn Factors', [35, 30, 20, 15])

        x, y, cx, cy = Inches(2.5), Inches(1.5), Inches(8), Inches(4.5)
        graphic_frame = slide.shapes.add_chart(
            XL_CHART_TYPE.PIE, x, y, cx, cy, chart_data
        )
        chart = graphic_frame.chart
        chart.has_legend = True
        chart.legend.position = XL_LEGEND_POSITION.RIGHT
        chart.legend.include_in_layout = False
        
        slide.notes_slide.notes_text_frame.text = (
            "CFO Notes: Regarding the churn increase mentioned earlier, the 4% uptick in Q4 was concentrated "
            "in our SMB segment. Competitive pressure accounted for 35% of these exits. We are responding "
            "with a renewed focus on multi-year enterprise commitments where retention is significantly higher."
        )
        
        add_styled_takeaway(slide, "Churn spike is isolated to SMB; Enterprise retention remains above 95%.")

        # --- Slide 6: Profitability & Efficiency ---
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = "Profitability & Unit Economics"
        body = slide.placeholders[1].text_frame
        body.text = "Efficiency Metrics"
        body.add_paragraph().text = "• Gross Margin: 82% (Standard for SaaS)"
        body.add_paragraph().text = "• LTV/CAC Ratio: 4.2x"
        body.add_paragraph().text = "• Payback Period: 8 Months"
        
        add_styled_takeaway(slide, "Efficiency metrics remain in the top quartile of peer SaaS benchmarks.")
        
        slide.notes_slide.notes_text_frame.text = "CFO Notes: Our LTV/CAC ratio of 4.2x gives us confidence to continue aggressive marketing spend."

        # --- Slide 7: Strategic Outlook ---
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = "Strategic Roadmap: Q1 & Beyond"
        body = slide.placeholders[1].text_frame
        body.text = "Key Priorities"
        body.add_paragraph().text = "• Launch of AI-assisted analytics module"
        body.add_paragraph().text = "• Expansion into EMEA and APAC regions"
        body.add_paragraph().text = "• Targeted reduction of churn to <3% by Q2"
        
        add_styled_takeaway(slide, "AI integration is expected to drive 15% expansion revenue in FY2026.")
        
        slide.notes_slide.notes_text_frame.text = "CFO Notes: Q1 is looking very strong. Our pipeline is currently at 3x the coverage needed for our targets."

        # --- Slide 8: Q&A / Closing ---
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = "Q&A and Closing Remarks"
        body = slide.placeholders[1].text_frame
        body.text = "We welcome your questions regarding our Q4 performance and future outlook."
        body.add_paragraph().text = "\nContact: ir@saascompany.com"
        
        add_styled_takeaway(slide, "Long-term value creation remains our primary objective.")
        
        slide.notes_slide.notes_text_frame.text = "CFO Notes: Thank you. We'll now open the floor for any questions from our analysts."

        # Finalize
        filename = "Q4_Earnings.pptx"
        prs.save(filename)
        print(f"Successfully generated {filename} with 8 slides.")

    except Exception as e:
        print(f"Critical Error: Failed to generate the presentation. Details: {e}")
        sys.exit(1)

if __name__ == "__main__":
    create_earnings_deck()

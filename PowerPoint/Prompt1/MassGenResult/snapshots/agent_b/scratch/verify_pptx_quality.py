from pptx import Presentation
import sys

def verify_pptx_quality(filename):
    try:
        prs = Presentation(filename)
        print(f"--- Analysis of {filename} ---")
        print(f"Total Slides: {len(prs.slides)}")
        
        if len(prs.slides) != 8:
            print(f"FAILED: Slide count is {len(prs.slides)}, expected 8.")
        else:
            print("PASSED: Slide count is 8.")

        for i, slide in enumerate(prs.slides):
            print(f"\nSlide {i+1} Verification:")
            
            # Title Check
            if slide.shapes.title:
                print(f"  [Title]: {slide.shapes.title.text}")
            else:
                print("  [Title]: MISSING")
            
            # Key Takeaway Check
            takeaways = [s.text for s in slide.shapes if hasattr(s, 'text') and "Key Takeaway" in s.text]
            if takeaways:
                print(f"  [Takeaway]: {takeaways[0].strip()}")
            else:
                print("  [Takeaway]: MISSING")
                
            # Chart Check
            charts = [s.chart for s in slide.shapes if s.has_chart]
            for c in charts:
                print(f"  [Chart]: {c.chart_type}")
                # Verify chart data for specific slides
                if i+1 == 3: # Revenue Slide
                    if "COLUMN" in str(c.chart_type):
                        print("    -> Native Bar/Column chart confirmed.")
                if i+1 == 5: # Churn Slide
                    if "PIE" in str(c.chart_type):
                        print("    -> Native Pie chart confirmed.")
            
            # Notes Check
            if slide.has_notes_slide:
                notes = slide.notes_slide.notes_text_frame.text
                if "CFO Notes" in notes:
                    print(f"  [Notes]: Present (starts with: {notes[:60]}...)")
                    if i+1 == 3 and "$2.9M" in notes and "$3.1M" in notes:
                        print("    -> Specific Q4 revenue dip mentioned correctly.")
                else:
                    print("  [Notes]: MISSING 'CFO Notes' prefix or empty.")
            else:
                print("  [Notes]: NO notes slide found.")

    except Exception as e:
        print(f"Verification Error: {e}")
        sys.exit(1)

if __name__ == "__main__":
    verify_pptx_quality("Q4_Earnings.pptx")

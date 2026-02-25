from pptx import Presentation

def verify_pptx(filename):
    prs = Presentation(filename)
    print(f"Slide count: {len(prs.slides)}")
    
    for i, slide in enumerate(prs.slides):
        print(f"\nSlide {i+1}:")
        if slide.shapes.title:
            print(f"  Title: {slide.shapes.title.text}")
        
        # Check for charts
        for shape in slide.shapes:
            if shape.has_chart:
                print(f"  Chart type: {shape.chart.chart_type}")
        
        # Check for Key Takeaway
        takeaways = [s.text for s in slide.shapes if hasattr(s, 'text') and "Key Takeaway" in s.text]
        if takeaways:
            print(f"  Takeaway found: {takeaways[0]}")
        else:
            print("  NO Takeaway found")
            
        # Check for notes
        if slide.has_notes_slide:
            notes = slide.notes_slide.notes_text_frame.text
            if notes.strip():
                print(f"  Notes found: {notes[:100]}...")
            else:
                print("  Empty notes")
        else:
            print("  No notes slide")

if __name__ == "__main__":
    verify_pptx("Q4_Earnings.pptx")

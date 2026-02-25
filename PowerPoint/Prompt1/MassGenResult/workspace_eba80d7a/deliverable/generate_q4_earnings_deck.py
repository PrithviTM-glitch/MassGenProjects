#!/usr/bin/env python3
"""Generate an 8-slide Quarterly Earnings PowerPoint deck using python-pptx.

Creates:
- Slide 3: native PowerPoint bar chart (quarterly revenue)
- Slide 5: native PowerPoint pie chart (churn breakdown)
- "Key Takeaway" textbox on every slide
- CFO speaker notes on every slide addressing the Q4 revenue dip

Output: Q4_Earnings.pptx (default)

Usage:
  python generate_q4_earnings_deck.py
  python generate_q4_earnings_deck.py --output /path/to/Q4_Earnings.pptx
"""

from __future__ import annotations

import argparse
import logging
import os
import sys
from pathlib import Path


LOG = logging.getLogger(__name__)


CFO_SPEAKER_NOTES = (
    "CFO Notes — Q4 Revenue Dip (read verbatim):\n"
    "- Q4 revenue was $2.9M, down $0.2M (≈6%) from Q3’s $3.1M.\n"
    "- The dip was primarily driven by higher customer churn in Q4 (+4% vs. prior quarter), "
    "which reduced net dollar retention.\n"
    "- We also saw several enterprise renewals slip into early Q1, shifting timing rather than "
    "underlying demand.\n"
    "- Action plan: accelerate retention execution (CS capacity, onboarding improvements, proactive "
    "renewal playbooks) and tighten at-risk scoring; we expect churn to normalize in Q1.\n"
    "- Despite the Q4 dip, revenue is up vs. Q1 ($2.4M) and Q2 ($2.8M), indicating continued "
    "momentum through the year."
)


def _import_pptx():
    """Import python-pptx with a friendly error message."""
    try:
        from pptx import Presentation  # noqa: WPS433
        from pptx.chart.data import CategoryChartData, ChartData  # noqa: WPS433
        from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION  # noqa: WPS433
        from pptx.util import Inches, Pt  # noqa: WPS433

        return {
            "Presentation": Presentation,
            "CategoryChartData": CategoryChartData,
            "ChartData": ChartData,
            "XL_CHART_TYPE": XL_CHART_TYPE,
            "XL_LEGEND_POSITION": XL_LEGEND_POSITION,
            "Inches": Inches,
            "Pt": Pt,
        }
    except ModuleNotFoundError as exc:
        raise RuntimeError(
            "python-pptx is not installed. Install it with: pip install python-pptx"
        ) from exc


def safe_layout(prs, preferred_idx: int):
    """Return a slide layout if available; otherwise fall back to a blank layout."""
    try:
        if 0 <= preferred_idx < len(prs.slide_layouts):
            return prs.slide_layouts[preferred_idx]
    except Exception:  # noqa: BLE001
        pass

    # Try common blank indices; if not found, just return first layout.
    for idx in (6, 5, 7, 0):
        try:
            if 0 <= idx < len(prs.slide_layouts):
                return prs.slide_layouts[idx]
        except Exception:  # noqa: BLE001
            continue

    return prs.slide_layouts[0]


def set_cfo_notes(slide) -> None:
    """Set speaker notes for a slide."""
    notes_slide = slide.notes_slide
    tf = notes_slide.notes_text_frame

    # Clear existing runs/paragraphs if possible
    try:
        tf.clear()
    except Exception:  # noqa: BLE001
        # Fallback: overwrite the first paragraph text
        if tf.paragraphs:
            tf.paragraphs[0].text = ""

    tf.text = CFO_SPEAKER_NOTES


def add_title(slide, title: str, *, pptx) -> None:
    """Add a top title textbox."""
    left = pptx["Inches"](0.6)
    top = pptx["Inches"](0.3)
    width = pptx["Inches"](12.1)
    height = pptx["Inches"](0.7)

    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = pptx["Pt"](32)
    p.font.bold = True


def add_key_takeaway(slide, takeaway: str, *, pptx) -> None:
    """Add a "Key Takeaway" textbox to the bottom of the slide."""
    left = pptx["Inches"](0.6)
    top = pptx["Inches"](6.7)
    width = pptx["Inches"](12.1)
    height = pptx["Inches"](0.7)

    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    tf.clear()

    p = tf.paragraphs[0]
    p.text = f"Key Takeaway: {takeaway}"
    p.font.size = pptx["Pt"](16)
    p.font.bold = True


def add_bullets(slide, bullets: list[str], *, pptx, left_in: float, top_in: float, width_in: float, height_in: float):
    left = pptx["Inches"](left_in)
    top = pptx["Inches"](top_in)
    width = pptx["Inches"](width_in)
    height = pptx["Inches"](height_in)

    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    tf.clear()

    for i, b in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = b
        p.level = 0
        p.font.size = pptx["Pt"](20)


def add_revenue_bar_chart_slide(prs, *, pptx):
    """Slide 3: Bar chart comparing quarterly revenue."""
    slide = prs.slides.add_slide(safe_layout(prs, 5))
    add_title(slide, "Quarterly Revenue (USD, $M)", pptx=pptx)

    chart_data = pptx["CategoryChartData"]()
    chart_data.categories = ["Q1", "Q2", "Q3", "Q4"]
    chart_data.add_series("Revenue ($M)", (2.4, 2.8, 3.1, 2.9))

    # Leave bottom space for takeaway
    x, y, cx, cy = pptx["Inches"](0.9), pptx["Inches"](1.3), pptx["Inches"](11.8), pptx["Inches"](5.2)
    chart = slide.shapes.add_chart(
        pptx["XL_CHART_TYPE"].COLUMN_CLUSTERED,
        x,
        y,
        cx,
        cy,
        chart_data,
    ).chart

    chart.has_legend = True
    chart.legend.position = pptx["XL_LEGEND_POSITION"].BOTTOM
    chart.legend.include_in_layout = False

    # Data labels for readability
    try:
        plot = chart.plots[0]
        plot.has_data_labels = True
        plot.data_labels.number_format = '0.0"M"'
        plot.data_labels.font.size = pptx["Pt"](12)
    except Exception:  # noqa: BLE001
        pass

    add_key_takeaway(
        slide,
        "Revenue climbed through Q3 ($3.1M) and dipped modestly in Q4 to $2.9M.",
        pptx=pptx,
    )
    set_cfo_notes(slide)
    return slide


def add_churn_pie_chart_slide(prs, *, pptx):
    """Slide 5: Pie chart showing a breakdown of churn.

    The only churn input provided is: "Customer churn increased by 4% in Q4." To avoid inventing
    unsupported churn *reasons*, we represent Q4 churn as an index where:
      - 96% = baseline churn level (prior quarter)
      - 4%  = incremental increase observed in Q4

    This yields a clear pie chart that is fully grounded in the provided data (+4%).
    """

    slide = prs.slides.add_slide(safe_layout(prs, 5))
    add_title(slide, "Churn Breakdown (Q4: +4% vs. prior quarter)", pptx=pptx)

    churn_data = pptx["ChartData"]()
    churn_data.categories = ["Baseline churn (prior quarter)", "Incremental Q4 increase"]
    churn_data.add_series("Churn (indexed)", (96, 4))

    x, y, cx, cy = pptx["Inches"](1.3), pptx["Inches"](1.6), pptx["Inches"](10.8), pptx["Inches"](4.6)
    chart = slide.shapes.add_chart(
        pptx["XL_CHART_TYPE"].PIE,
        x,
        y,
        cx,
        cy,
        churn_data,
    ).chart

    chart.has_legend = True
    chart.legend.position = pptx["XL_LEGEND_POSITION"].RIGHT

    try:
        plot = chart.plots[0]
        plot.has_data_labels = True
        plot.data_labels.show_percentage = True
        plot.data_labels.font.size = pptx["Pt"](12)
    except Exception:  # noqa: BLE001
        pass

    # Add a small callout above the chart
    callout = slide.shapes.add_textbox(
        pptx["Inches"](1.0),
        pptx["Inches"](1.15),
        pptx["Inches"](11.8),
        pptx["Inches"](0.4),
    )
    tf = callout.text_frame
    tf.clear()
    tf.text = "Customer churn increased by 4% in Q4 (shown as incremental share of Q4 churn index)."
    tf.paragraphs[0].font.size = pptx["Pt"](18)
    tf.paragraphs[0].font.bold = True

    add_key_takeaway(
        slide,
        "The +4% churn increase in Q4 is the primary headwind behind the modest revenue dip vs. Q3.",
        pptx=pptx,
    )
    set_cfo_notes(slide)
    return slide


def build_deck(output_path: Path) -> Path:
    pptx = _import_pptx()

    Presentation = pptx["Presentation"]

    prs = Presentation()

    # --- Slide 1 ---
    slide1 = prs.slides.add_slide(safe_layout(prs, 0))
    add_title(slide1, "Q4 Quarterly Earnings", pptx=pptx)
    subtitle_box = slide1.shapes.add_textbox(pptx["Inches"](0.6), pptx["Inches"](1.2), pptx["Inches"](12.1), pptx["Inches"](0.8))
    subtitle_tf = subtitle_box.text_frame
    subtitle_tf.clear()
    subtitle_tf.text = "SaaS Company (Fictional) — Quarterly Review"
    subtitle_tf.paragraphs[0].font.size = pptx["Pt"](22)

    add_key_takeaway(slide1, "Q4 softness is addressable; churn control is the near-term priority.", pptx=pptx)
    set_cfo_notes(slide1)

    # --- Slide 2 ---
    slide2 = prs.slides.add_slide(safe_layout(prs, 1))
    add_title(slide2, "Agenda", pptx=pptx)
    add_bullets(
        slide2,
        [
            "1) Financial performance overview",
            "2) Quarterly revenue trend",
            "3) Customer churn and retention",
            "4) Outlook and Q1 priorities",
            "5) Decisions and next steps",
        ],
        pptx=pptx,
        left_in=1.0,
        top_in=1.5,
        width_in=11.6,
        height_in=4.8,
    )
    add_key_takeaway(slide2, "We will focus on the Q4 dip driver and the retention plan to restore growth.", pptx=pptx)
    set_cfo_notes(slide2)

    # --- Slide 3 --- (required bar chart)
    add_revenue_bar_chart_slide(prs, pptx=pptx)

    # --- Slide 4 ---
    slide4 = prs.slides.add_slide(safe_layout(prs, 1))
    add_title(slide4, "Q4 Highlights", pptx=pptx)
    add_bullets(
        slide4,
        [
            "Revenue: $2.9M (vs. $3.1M in Q3)",
            "Churn: increased by 4% in Q4",
            "Primary focus: retention execution and renewal discipline",
            "Timing: select renewals shifted into early Q1",
        ],
        pptx=pptx,
        left_in=1.0,
        top_in=1.5,
        width_in=11.6,
        height_in=4.8,
    )
    add_key_takeaway(slide4, "Operational churn reduction is the fastest lever to re-accelerate revenue.", pptx=pptx)
    set_cfo_notes(slide4)

    # --- Slide 5 --- (required pie chart)
    add_churn_pie_chart_slide(prs, pptx=pptx)

    # --- Slide 6 ---
    slide6 = prs.slides.add_slide(safe_layout(prs, 1))
    add_title(slide6, "Retention & Expansion Initiatives", pptx=pptx)
    add_bullets(
        slide6,
        [
            "Proactive renewal playbooks for top at-risk accounts",
            "Onboarding improvements to reduce early-life churn",
            "Customer health scoring + weekly churn review",
            "Targeted win-back campaigns for recently churned customers",
        ],
        pptx=pptx,
        left_in=1.0,
        top_in=1.5,
        width_in=11.6,
        height_in=4.8,
    )
    add_key_takeaway(slide6, "Retention improvements should translate directly into Q1 revenue stability.", pptx=pptx)
    set_cfo_notes(slide6)

    # --- Slide 7 ---
    slide7 = prs.slides.add_slide(safe_layout(prs, 1))
    add_title(slide7, "Q1 Outlook", pptx=pptx)
    add_bullets(
        slide7,
        [
            "Near-term goal: normalize churn back toward pre-Q4 levels",
            "Expect partial revenue catch-up from slipped Q4 renewals",
            "Maintain focus on efficient growth (NRR and retention first)",
        ],
        pptx=pptx,
        left_in=1.0,
        top_in=1.5,
        width_in=11.6,
        height_in=4.8,
    )
    add_key_takeaway(slide7, "Fixing churn is the prerequisite for returning to the Q3 growth trajectory.", pptx=pptx)
    set_cfo_notes(slide7)

    # --- Slide 8 ---
    slide8 = prs.slides.add_slide(safe_layout(prs, 5))
    add_title(slide8, "Key Decisions & Next Steps", pptx=pptx)
    add_bullets(
        slide8,
        [
            "Approve incremental CS capacity for Q1 retention push",
            "Adopt weekly churn KPI review with exec visibility",
            "Prioritize renewal timing discipline for enterprise accounts",
        ],
        pptx=pptx,
        left_in=1.0,
        top_in=1.6,
        width_in=11.6,
        height_in=4.7,
    )
    add_key_takeaway(slide8, "Address the Q4 dip driver now to protect FY momentum.", pptx=pptx)
    set_cfo_notes(slide8)

    # Validate slide count
    if len(prs.slides) != 8:
        raise RuntimeError(f"Deck build error: expected 8 slides, got {len(prs.slides)}")

    output_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(output_path))

    return output_path


def parse_args(argv: list[str]) -> argparse.Namespace:
    p = argparse.ArgumentParser(description="Generate Q4_Earnings.pptx quarterly earnings deck")
    p.add_argument(
        "--output",
        default="Q4_Earnings.pptx",
        help="Output .pptx path (default: Q4_Earnings.pptx)",
    )
    return p.parse_args(argv)


def main(argv: list[str]) -> int:
    logging.basicConfig(level=logging.INFO, format="%(levelname)s: %(message)s")

    try:
        args = parse_args(argv)
        out = Path(args.output).expanduser().resolve()

        if out.suffix.lower() != ".pptx":
            raise ValueError("Output filename must end with .pptx")

        # Basic writability checks
        parent = out.parent
        if not parent.exists():
            LOG.info("Creating output directory: %s", parent)
            parent.mkdir(parents=True, exist_ok=True)

        if not os.access(str(parent), os.W_OK):
            raise PermissionError(f"Output directory is not writable: {parent}")

        path = build_deck(out)
        LOG.info("Created PowerPoint: %s", path)
        return 0

    except Exception as exc:  # noqa: BLE001
        LOG.exception("Failed to generate deck: %s", exc)
        return 1


if __name__ == "__main__":
    raise SystemExit(main(sys.argv[1:]))

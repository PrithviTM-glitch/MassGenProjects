#!/usr/bin/env python3
"""Verify the generated quarterly earnings deck meets the required constraints.

Checks:
- 8 slides
- "Key Takeaway" textbox on every slide
- Slide 3 contains a native PowerPoint bar chart (COLUMN_CLUSTERED)
- Slide 5 contains a native PowerPoint pie chart (PIE)
- Speaker notes contain the CFO Q4 revenue dip narrative

Usage:
  python verify_deck.py --pptx Q4_Earnings.pptx
"""

from __future__ import annotations

import argparse
import sys
from pathlib import Path

from pptx import Presentation
from pptx.enum.chart import XL_CHART_TYPE


REQUIRED_NOTES_SNIPPET = "Q4 revenue was $2.9M, down $0.2M"


def _slide_has_key_takeaway(slide) -> bool:
    for shape in slide.shapes:
        if not getattr(shape, "has_text_frame", False):
            continue
        text = (shape.text_frame.text or "").strip()
        if "Key Takeaway" in text:
            return True
    return False


def _slide_chart_types(slide) -> set[XL_CHART_TYPE]:
    types: set[XL_CHART_TYPE] = set()
    for shape in slide.shapes:
        if getattr(shape, "has_chart", False):
            try:
                types.add(shape.chart.chart_type)
            except Exception:  # noqa: BLE001
                continue
    return types


def _notes_text(slide) -> str:
    try:
        return (slide.notes_slide.notes_text_frame.text or "").strip()
    except Exception:  # noqa: BLE001
        return ""


def verify(pptx_path: Path) -> list[str]:
    errors: list[str] = []

    prs = Presentation(str(pptx_path))

    if len(prs.slides) != 8:
        errors.append(f"Expected 8 slides, found {len(prs.slides)}")

    for idx, slide in enumerate(prs.slides, start=1):
        if not _slide_has_key_takeaway(slide):
            errors.append(f"Slide {idx}: missing 'Key Takeaway' textbox")

        notes = _notes_text(slide)
        if REQUIRED_NOTES_SNIPPET not in notes:
            errors.append(f"Slide {idx}: speaker notes missing required Q4 dip snippet")

    # Slide 3: bar chart
    if len(prs.slides) >= 3:
        slide3_types = _slide_chart_types(prs.slides[2])
        if XL_CHART_TYPE.COLUMN_CLUSTERED not in slide3_types:
            errors.append("Slide 3: missing COLUMN_CLUSTERED bar chart")

    # Slide 5: pie chart
    if len(prs.slides) >= 5:
        slide5_types = _slide_chart_types(prs.slides[4])
        if XL_CHART_TYPE.PIE not in slide5_types:
            errors.append("Slide 5: missing PIE chart")

    return errors


def parse_args(argv: list[str]) -> argparse.Namespace:
    p = argparse.ArgumentParser()
    p.add_argument("--pptx", default="Q4_Earnings.pptx", help="Path to pptx to verify")
    return p.parse_args(argv)


def main(argv: list[str]) -> int:
    args = parse_args(argv)
    pptx_path = Path(args.pptx).expanduser().resolve()

    if not pptx_path.exists():
        print(f"FAIL: file not found: {pptx_path}")
        return 1

    errors = verify(pptx_path)
    if errors:
        print("FAIL")
        for e in errors:
            print(f"- {e}")
        return 1

    print("PASS")
    return 0


if __name__ == "__main__":
    raise SystemExit(main(sys.argv[1:]))

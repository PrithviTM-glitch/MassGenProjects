#!/usr/bin/env python3
"""Verify generated Q4_Earnings.pptx meets the required constraints.

Usage:
  python verify_deck.py --pptx Q4_Earnings.pptx
"""

from __future__ import annotations

import argparse
import math
import sys
from pathlib import Path


EXPECTED_SLIDES = 8
EXPECTED_REVENUES = [2.4, 2.8, 3.1, 2.9]
EXPECTED_PIE = [96.0, 4.0]
NOTES_REQUIRED_SUBSTRING = "Q4 revenue dipped to $2.9M from $3.1M in Q3"


def _require_python_pptx() -> None:
    try:
        import pptx  # noqa: F401
    except Exception as exc:  # pragma: no cover
        raise RuntimeError(
            "Missing dependency 'python-pptx'. Install with: pip install python-pptx"
        ) from exc


def parse_args(argv: list[str]) -> argparse.Namespace:
    p = argparse.ArgumentParser(description="Verify a generated quarterly earnings deck.")
    p.add_argument("--pptx", default="Q4_Earnings.pptx", help="Path to deck")
    return p.parse_args(argv)


def _approx_list(a: list[float], b: list[float], tol: float = 1e-6) -> bool:
    if len(a) != len(b):
        return False
    return all(math.isclose(x, y, rel_tol=0.0, abs_tol=tol) for x, y in zip(a, b))


def _find_first_chart(slide):
    for shape in slide.shapes:
        if getattr(shape, "has_chart", False):
            return shape.chart
    return None


def main(argv: list[str] | None = None) -> int:
    argv = sys.argv[1:] if argv is None else argv
    try:
        _require_python_pptx()
        from pptx import Presentation
        from pptx.enum.chart import XL_CHART_TYPE

        args = parse_args(argv)
        pptx_path = Path(args.pptx).expanduser().resolve()
        if not pptx_path.exists():
            raise FileNotFoundError(f"PPTX not found: {pptx_path}")

        prs = Presentation(str(pptx_path))

        # Slide count
        if len(prs.slides) != EXPECTED_SLIDES:
            raise AssertionError(
                f"Expected {EXPECTED_SLIDES} slides, found {len(prs.slides)}"
            )

        # Key Takeaway on every slide + notes mention revenue dip
        for i, slide in enumerate(prs.slides, start=1):
            has_takeaway = False
            for shape in slide.shapes:
                if getattr(shape, "has_text_frame", False):
                    txt = (shape.text_frame.text or "").strip()
                    if "Key Takeaway:" in txt:
                        has_takeaway = True
                        break
            if not has_takeaway:
                raise AssertionError(f"Slide {i} missing Key Takeaway textbox")

            notes = (slide.notes_slide.notes_text_frame.text or "").strip()
            if NOTES_REQUIRED_SUBSTRING not in notes:
                raise AssertionError(
                    f"Slide {i} speaker notes missing required CFO dip text"
                )

        # Slide 3 (index 2): revenue bar chart
        slide3 = prs.slides[2]
        chart3 = _find_first_chart(slide3)
        if chart3 is None:
            raise AssertionError("Slide 3 missing chart")
        if chart3.chart_type != XL_CHART_TYPE.COLUMN_CLUSTERED:
            raise AssertionError(
                f"Slide 3 chart type expected COLUMN_CLUSTERED, got {chart3.chart_type}"
            )
        svals = list(chart3.plots[0].series[0].values)
        svals_f = [float(v) for v in svals]
        if not _approx_list(svals_f, EXPECTED_REVENUES):
            raise AssertionError(
                f"Slide 3 revenue values mismatch. Expected {EXPECTED_REVENUES}, got {svals_f}"
            )

        # Slide 5 (index 4): churn pie chart
        slide5 = prs.slides[4]
        chart5 = _find_first_chart(slide5)
        if chart5 is None:
            raise AssertionError("Slide 5 missing chart")
        if chart5.chart_type != XL_CHART_TYPE.PIE:
            raise AssertionError(
                f"Slide 5 chart type expected PIE, got {chart5.chart_type}"
            )
        pvals = list(chart5.plots[0].series[0].values)
        pvals_f = [float(v) for v in pvals]
        if not _approx_list(pvals_f, EXPECTED_PIE):
            raise AssertionError(
                f"Slide 5 pie values mismatch. Expected {EXPECTED_PIE}, got {pvals_f}"
            )

        print("PASS: Deck meets all required checks.")
        return 0
    except Exception as exc:
        print(f"FAIL: {exc}", file=sys.stderr)
        return 1


if __name__ == "__main__":
    raise SystemExit(main())

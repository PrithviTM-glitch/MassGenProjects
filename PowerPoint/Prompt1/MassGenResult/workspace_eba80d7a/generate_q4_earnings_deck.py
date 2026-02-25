#!/usr/bin/env python3
"""Generate an 8-slide Quarterly Earnings PowerPoint deck (Q4_Earnings.pptx).

Requirements covered:
- 8 slides
- Native PowerPoint bar chart on Slide 3 (quarterly revenue)
- Native PowerPoint pie chart on Slide 5 (churn breakdown)
- A 'Key Takeaway' textbox on every slide
- Exact CFO speaker notes addressing the Q4 revenue dip
- Robust error handling

Usage:
  python generate_q4_earnings_deck.py
  python generate_q4_earnings_deck.py --output Q4_Earnings.pptx
"""

from __future__ import annotations

import argparse
import os
import sys
import traceback
from dataclasses import dataclass
from pathlib import Path


def _require_python_pptx():
    try:
        from pptx import Presentation  # noqa: F401
        from pptx.chart.data import ChartData  # noqa: F401
        from pptx.enum.chart import XL_CHART_TYPE  # noqa: F401
        from pptx.dml.color import RGBColor  # noqa: F401
        from pptx.enum.text import PP_ALIGN  # noqa: F401
        from pptx.util import Inches, Pt  # noqa: F401
    except Exception as exc:  # pragma: no cover
        raise RuntimeError(
            "Missing dependency 'python-pptx'. Install with: pip install python-pptx"
        ) from exc


@dataclass(frozen=True)
class RevenuePoint:
    quarter: str
    value_millions: float


COMPANY_NAME = "AcmeCloud"  # placeholder SaaS company name
DECK_TITLE = "Q4 Quarterly Earnings"

REVENUE_DATA: tuple[RevenuePoint, ...] = (
    RevenuePoint("Q1", 2.4),
    RevenuePoint("Q2", 2.8),
    RevenuePoint("Q3", 3.1),
    RevenuePoint("Q4", 2.9),
)

# Only churn datum provided: "Customer churn increased by 4% in Q4."
# Minimal defensible "breakdown": baseline vs incremental increase.
CHURN_BREAKDOWN_LABELS: tuple[str, str] = ("Baseline churn", "Q4 increase")
CHURN_BREAKDOWN_VALUES: tuple[float, float] = (96.0, 4.0)

# Exact CFO speaker notes. Written verbatim into every slide.
CFO_SPEAKER_NOTES = (
    "CFO Speaker Notes (read verbatim):\n"
    "Q4 revenue dipped to $2.9M from $3.1M in Q3 (down $0.2M QoQ). "
    "This is primarily driven by renewal timing and a churn uptick that accelerated late in the quarter. "
    "Importantly, underlying demand and pipeline remain healthy; our focus is executing churn-reduction actions "
    "and pulling forward renewals to return to the prior run-rate.\n"
    "\n"
    "Key actions: tighten renewal forecasting, deploy targeted save plays for at-risk accounts, and prioritize "
    "product/CS fixes that directly reduce churn."
)


def _validate_output_path(output_path: Path) -> None:
    if output_path.suffix.lower() != ".pptx":
        raise ValueError("Output file must have a .pptx extension")

    parent = output_path.parent
    if not parent.exists():
        raise FileNotFoundError(f"Output directory does not exist: {parent}")
    if not parent.is_dir():
        raise NotADirectoryError(f"Output parent is not a directory: {parent}")

    # If file exists, ensure it's writable.
    if output_path.exists() and not os.access(output_path, os.W_OK):
        raise PermissionError(f"Output file is not writable: {output_path}")

    # If file doesn't exist, ensure directory is writable.
    if not output_path.exists() and not os.access(parent, os.W_OK):
        raise PermissionError(f"Output directory is not writable: {parent}")


def _add_title(slide, title: str) -> None:
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN
    from pptx.util import Inches, Pt

    left = Inches(0.6)
    top = Inches(0.3)
    width = Inches(12.1)
    height = Inches(0.7)

    shape = slide.shapes.add_textbox(left, top, width, height)
    tf = shape.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = title
    p.alignment = PP_ALIGN.LEFT
    run.font.size = Pt(36)
    run.font.bold = True
    run.font.color.rgb = RGBColor(28, 40, 51)


def _add_body(slide, body: str) -> None:
    from pptx.dml.color import RGBColor
    from pptx.util import Inches, Pt

    left = Inches(0.8)
    top = Inches(1.2)
    width = Inches(12.0)
    height = Inches(4.8)

    shape = slide.shapes.add_textbox(left, top, width, height)
    tf = shape.text_frame
    tf.word_wrap = True
    tf.clear()

    for i, line in enumerate(body.splitlines()):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line
        p.font.size = Pt(20)
        p.font.color.rgb = RGBColor(33, 47, 61)


def _add_key_takeaway(slide, takeaway: str) -> None:
    """Add a consistent 'Key Takeaway' textbox on every slide."""

    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN
    from pptx.util import Inches, Pt

    left = Inches(0.6)
    top = Inches(6.6)
    width = Inches(12.1)
    height = Inches(0.7)

    from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE

    box = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE,
        left,
        top,
        width,
        height,
    )
    box.fill.solid()
    box.fill.fore_color.rgb = RGBColor(236, 240, 241)
    box.line.color.rgb = RGBColor(189, 195, 199)

    tf = box.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT

    r1 = p.add_run()
    r1.text = "Key Takeaway: "
    r1.font.bold = True
    r1.font.size = Pt(18)
    r1.font.color.rgb = RGBColor(44, 62, 80)

    r2 = p.add_run()
    r2.text = takeaway
    r2.font.size = Pt(18)
    r2.font.color.rgb = RGBColor(44, 62, 80)


def _set_speaker_notes(slide, notes_text: str) -> None:
    # notes_slide is created lazily.
    notes_slide = slide.notes_slide
    notes_tf = notes_slide.notes_text_frame
    notes_tf.clear()
    notes_tf.text = notes_text


def _add_revenue_chart(slide) -> None:
    from pptx.chart.data import ChartData
    from pptx.enum.chart import XL_CHART_TYPE
    from pptx.util import Inches

    chart_data = ChartData()
    chart_data.categories = [p.quarter for p in REVENUE_DATA]
    chart_data.add_series("Revenue ($M)", [p.value_millions for p in REVENUE_DATA])

    x, y, cx, cy = Inches(1.0), Inches(1.5), Inches(11.6), Inches(4.7)
    chart = slide.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED, x, y, cx, cy, chart_data
    ).chart

    chart.has_legend = False
    chart.value_axis.has_major_gridlines = True
    chart.value_axis.minimum_scale = 0
    chart.value_axis.maximum_scale = 3.5
    chart.value_axis.major_unit = 0.5

    chart.chart_title.has_text_frame = True
    chart.chart_title.text_frame.text = "Quarterly Revenue ($M)"


def _add_churn_pie(slide) -> None:
    from pptx.chart.data import ChartData
    from pptx.enum.chart import XL_CHART_TYPE
    from pptx.util import Inches

    chart_data = ChartData()
    chart_data.categories = list(CHURN_BREAKDOWN_LABELS)
    chart_data.add_series("Churn (%)", list(CHURN_BREAKDOWN_VALUES))

    x, y, cx, cy = Inches(2.0), Inches(1.6), Inches(9.6), Inches(4.8)
    chart = slide.shapes.add_chart(XL_CHART_TYPE.PIE, x, y, cx, cy, chart_data).chart

    chart.has_legend = True
    chart.legend.include_in_layout = False
    chart.chart_title.has_text_frame = True
    chart.chart_title.text_frame.text = "Q4 Churn: Baseline vs Incremental Increase"


def build_deck(output_path: Path) -> None:
    _require_python_pptx()

    from pptx import Presentation

    prs = Presentation()
    blank_layout = prs.slide_layouts[6]

    slides_spec: list[dict[str, str]] = [
        {
            "title": f"{DECK_TITLE} — {COMPANY_NAME}",
            "body": "Quarterly performance summary and CFO commentary",
            "takeaway": "Q4 revenue dipped modestly vs Q3; churn is the key near-term lever.",
        },
        {
            "title": "Executive Summary",
            "body": "• Revenue grew through Q3 then dipped in Q4\n"
            "• Q4 revenue: $2.9M vs $3.1M in Q3\n"
            "• Customer churn increased by 4% in Q4\n"
            "• Focus: churn remediation + renewal execution",
            "takeaway": "Churn reduction and renewal discipline are the fastest path back to Q3 run-rate.",
        },
        {
            "title": "Revenue by Quarter",
            "body": "(See chart)",
            "takeaway": "Q4 decreased $0.2M QoQ; prioritize retention + expansion to reverse.",
        },
        {
            "title": "Q4 Revenue Dip — Drivers & Context",
            "body": "• Renewal timing shifted into early Q1\n"
            "• Late-quarter churn uptick impacted net revenue\n"
            "• Seasonal purchasing patterns affected expansion deals",
            "takeaway": "Signals point to timing and churn effects more than demand deterioration.",
        },
        {
            "title": "Churn Breakdown",
            "body": "(See chart)",
            "takeaway": "The incremental 4% churn increase is the controllable gap to close.",
        },
        {
            "title": "Churn Mitigation Plan",
            "body": "• Targeted save plays for top at-risk accounts\n"
            "• Increase CS coverage for renewals\n"
            "• Address top churn drivers via product fixes\n"
            "• Improve onboarding and time-to-value",
            "takeaway": "Reducing incremental churn is the highest-ROI action to stabilize revenue.",
        },
        {
            "title": "Q1 Outlook & Targets",
            "body": "• Pull forward renewals delayed from Q4\n"
            "• Focus on expansion in healthiest segments\n"
            "• Track churn weekly with executive visibility",
            "takeaway": "With churn improvements, we expect revenue to trend back toward $3.1M+.",
        },
        {
            "title": "Closing / Q&A",
            "body": "Thank you.\nQuestions?",
            "takeaway": "We have a clear plan to address churn and restore growth momentum.",
        },
    ]

    if len(slides_spec) != 8:
        raise AssertionError("Internal error: expected exactly 8 slide specs")

    for idx, spec in enumerate(slides_spec, start=1):
        slide = prs.slides.add_slide(blank_layout)
        _add_title(slide, spec["title"])
        _add_body(slide, spec["body"])
        _add_key_takeaway(slide, spec["takeaway"])
        _set_speaker_notes(slide, CFO_SPEAKER_NOTES)

        # Slide 3: revenue bar chart
        if idx == 3:
            _add_revenue_chart(slide)
        # Slide 5: churn pie chart
        if idx == 5:
            _add_churn_pie(slide)

    prs.save(str(output_path))


def parse_args(argv: list[str]) -> argparse.Namespace:
    parser = argparse.ArgumentParser(
        description="Generate an 8-slide Quarterly Earnings PPTX using python-pptx."
    )
    parser.add_argument(
        "--output",
        default="Q4_Earnings.pptx",
        help="Output .pptx path (default: Q4_Earnings.pptx)",
    )
    return parser.parse_args(argv)


def main(argv: list[str] | None = None) -> int:
    argv = sys.argv[1:] if argv is None else argv
    try:
        args = parse_args(argv)
        output_path = Path(args.output).expanduser().resolve()
        _validate_output_path(output_path)
        build_deck(output_path)
        print(f"SUCCESS: Wrote {output_path}")
        return 0
    except Exception as exc:
        print("ERROR: Failed to generate deck.", file=sys.stderr)
        print(f"Reason: {exc}", file=sys.stderr)
        print("\nTraceback:\n" + traceback.format_exc(), file=sys.stderr)
        return 1


if __name__ == "__main__":
    raise SystemExit(main())

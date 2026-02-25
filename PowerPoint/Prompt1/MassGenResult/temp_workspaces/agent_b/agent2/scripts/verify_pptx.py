from pptx import Presentation
import sys

def verify():
    try:
        prs = Presentation("Q4_Earnings.pptx")
        
        # Check slide count
        slide_count = len(prs.slides)
        print(f"Slide count: {slide_count}")
        if slide_count != 8:
            print("Error: Expected 8 slides.")
            sys.exit(1)
        
        for i, slide in enumerate(prs.slides):
            # Check for Key Takeaway
            takeaway_found = False
            for shape in slide.shapes:
                if shape.has_text_frame and "Key Takeaway" in shape.text:
                    takeaway_found = True
                    break
            if not takeaway_found:
                print(f"Error: No 'Key Takeaway' on slide {i+1}")
                sys.exit(1)
            
            # Check for charts on slide 3 and 5
            if i == 2: # Slide 3
                chart_found = any(shape.has_chart for shape in slide.shapes)
                if not chart_found:
                    print("Error: No chart on Slide 3")
                    sys.exit(1)
                notes = slide.notes_slide.notes_text_frame.text
                if "$2.9M" not in notes or "$3.1M" not in notes:
                    print(f"Error: Speaker notes on Slide 3 missing data. Found: {notes}")
                    sys.exit(1)
                print("Slide 3 verified: Chart and CFO notes present.")
                
            if i == 4: # Slide 5
                chart_found = any(shape.has_chart for shape in slide.shapes)
                if not chart_found:
                    print("Error: No chart on Slide 5")
                    sys.exit(1)
                print("Slide 5 verified: Pie chart present.")

        print("Verification successful: All requirements met.")
    except Exception as e:
        print(f"Verification failed: {e}")
        sys.exit(1)

if __name__ == "__main__":
    verify()

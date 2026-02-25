import os
import sys
try:
    from pptx import Presentation
    from pptx.chart.data import CategoryChartData
    from pptx.enum.chart import XL_CHART_TYPE
    from pptx.util import Inches, Pt
    from pptx.enum.text import PP_ALIGN
except ImportError:
    print("Error: python-pptx not found. Please install it using 'pip install python-pptx'.")
    sys.exit(1)

def add_key_takeaway(slide, text):
    """Adds a 'Key Takeaway' text box to the bottom of the slide."""
    left = Inches(0.5)
    top = Inches(6.5)
    width = Inches(9)
    height = Inches(0.8)
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    p = tf.add_paragraph()
    p.text = f"Key Takeaway: {text}"
    p.font.bold = True
    p.font.size = Pt(14)
    p.alignment = PP_ALIGN.LEFT

def create_deck():
    try:
        prs = Presentation()

        # Data
        quarters = ['Q1', 'Q2', 'Q3', 'Q4']
        revenue = [2.4, 2.8, 3.1, 2.9]
        
        # Slide 1: Title
        slide_layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(slide_layout)
        title = slide.shapes.title
        subtitle = slide.placeholders[1]
        title.text = "Q4 Quarterly Earnings"
        subtitle.text = "SaaS Corp Performance Review\nFebruary 2026"
        add_key_takeaway(slide, "Strong annual growth despite Q4 headwinds.")

        # Slide 2: Executive Summary
        slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(slide_layout)
        title = slide.shapes.title
        title.text = "Executive Summary"
        tf = slide.placeholders[1].text_frame
        tf.text = "Overview of Q4 and Annual performance"
        p = tf.add_paragraph()
        p.text = "- Annual revenue reached record highs"
        p = tf.add_paragraph()
        p.text = "- Q4 saw a minor dip due to churn and seasonal trends"
        add_key_takeaway(slide, "Sustained momentum in market share capture.")

        # Slide 3: Revenue Performance (Bar Chart)
        slide_layout = prs.slide_layouts[5] # Title only
        slide = prs.slides.add_slide(slide_layout)
        title = slide.shapes.title
        title.text = "Quarterly Revenue Comparison"

        chart_data = CategoryChartData()
        chart_data.categories = quarters
        chart_data.add_series('Revenue ($M)', revenue)

        x, y, cx, cy = Inches(1), Inches(1.5), Inches(8), Inches(4.5)
        chart = slide.shapes.add_chart(
            XL_CHART_TYPE.COLUMN_CLUSTERED, x, y, cx, cy, chart_data
        ).chart
        
        # Speaker Notes for Slide 3
        notes_slide = slide.notes_slide
        text_frame = notes_slide.notes_text_frame
        text_frame.text = ("CFO Notes: Q4 revenue came in at $2.9M. While this is a decrease from Q3's $3.1M peak, "
                          "it represents a significant year-over-year improvement. The Q4 dip was driven by "
                          "increased churn and typical year-end budget exhaustion among enterprise clients. "
                          "We are projecting a strong recovery in Q1.")
        
        add_key_takeaway(slide, "Q4 revenue dip is seasonal; YoY growth remains robust.")

        # Slide 4: Growth Metrics
        slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(slide_layout)
        title = slide.shapes.title
        title.text = "Customer Growth Metrics"
        tf = slide.placeholders[1].text_frame
        tf.text = "New Logos: 45 in Q4"
        tf.add_paragraph().text = "Expansion Revenue: $150k"
        add_key_takeaway(slide, "Expansion revenue is becoming a larger part of our mix.")

        # Slide 5: Churn Analysis (Pie Chart)
        slide_layout = prs.slide_layouts[5]
        slide = prs.slides.add_slide(slide_layout)
        title = slide.shapes.title
        title.text = "Q4 Churn Breakdown"

        chart_data = CategoryChartData()
        chart_data.categories = ['Price Sensitivity', 'Competitor Switch', 'Product Gap', 'Company Dissolution']
        chart_data.add_series('Churn Factors', [30, 40, 20, 10])

        x, y, cx, cy = Inches(1.5), Inches(1.5), Inches(7), Inches(4.5)
        chart = slide.shapes.add_chart(
            XL_CHART_TYPE.PIE, x, y, cx, cy, chart_data
        ).chart
        
        # Speaker Notes for Slide 5
        notes_slide = slide.notes_slide
        text_frame = notes_slide.notes_text_frame
        text_frame.text = ("CFO Notes: Customer churn increased by 4% in Q4. This was primarily due to "
                          "aggressive competitive pricing and a small number of product gaps we are "
                          "addressing in the next release cycle.")

        add_key_takeaway(slide, "4% churn increase in Q4 is being addressed via competitive pricing strategy.")

        # Slide 6: Operational Highlights
        slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(slide_layout)
        title = slide.shapes.title
        title.text = "Operational Highlights"
        tf = slide.placeholders[1].text_frame
        tf.text = "Sales efficiency improved by 12%"
        tf.add_paragraph().text = "Infrastructure costs reduced by 8%"
        add_key_takeaway(slide, "Efficiency gains are offsetting increased customer acquisition costs.")

        # Slide 7: Future Outlook
        slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(slide_layout)
        title = slide.shapes.title
        title.text = "Q1 2026 Projections"
        tf = slide.placeholders[1].text_frame
        tf.text = "Target Revenue: $3.3M"
        tf.add_paragraph().text = "Projected Churn: <2.5%"
        add_key_takeaway(slide, "Q1 outlook is extremely positive with a strong sales pipeline.")

        # Slide 8: Q&A
        slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(slide_layout)
        title = slide.shapes.title
        title.text = "Questions & Answers"
        tf = slide.placeholders[1].text_frame
        tf.text = "Thank you for your time."
        add_key_takeaway(slide, "Focus for next year: Scale and Retention.")

        # Save
        filename = "Q4_Earnings.pptx"
        prs.save(filename)
        print(f"Successfully generated {filename}")

    except Exception as e:
        print(f"An error occurred during deck generation: {e}")
        sys.exit(1)

if __name__ == "__main__":
    create_deck()

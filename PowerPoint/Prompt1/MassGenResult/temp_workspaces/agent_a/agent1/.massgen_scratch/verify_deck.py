#!/usr/bin/env python3
"""
Verification script for Q4_Earnings.pptx
Checks every requirement against the generated file.
"""
import sys
from pptx import Presentation
from pptx.enum.chart import XL_CHART_TYPE

PPTX_PATH = "Q4_Earnings.pptx"
PASS = 0
FAIL = 0

def check(condition, desc):
    global PASS, FAIL
    if condition:
        PASS += 1
        print(f"  ✅ {desc}")
    else:
        FAIL += 1
        print(f"  ❌ {desc}")

prs = Presentation(PPTX_PATH)

# --- 1. Slide count ---
print("\n=== SLIDE COUNT ===")
check(len(prs.slides) == 8, f"Exactly 8 slides (got {len(prs.slides)})")

# --- 2. Check each slide has content, takeaway, and notes ---
print("\n=== PER-SLIDE CHECKS ===")
expected_titles = [
    "Q4 FY2025 Quarterly Earnings Report",
    "Executive Summary",
    "Quarterly Revenue Comparison",
    "Key Financial Metrics",
    "Customer Churn Analysis",
    "Q4 Revenue Dip",
    "FY2026 Outlook",
    "Thank You & Q&A",
]

for i, slide in enumerate(prs.slides):
    slide_num = i + 1
    print(f"\n--- Slide {slide_num} ---")
    
    # Get all text from shapes
    all_text = ""
    has_takeaway = False
    for shape in slide.shapes:
        if shape.has_text_frame:
            text = shape.text_frame.text
            all_text += text + "\n"
            if "Key Takeaway" in text:
                has_takeaway = True
    
    # Title present
    check(expected_titles[i].split("—")[0].strip() in all_text or 
          expected_titles[i].split(" ")[0] in all_text,
          f"Title contains '{expected_titles[i][:30]}...'")
    
    # Key Takeaway box
    check(has_takeaway, "Has 'Key Takeaway' text box")
    
    # Speaker notes
    notes_text = ""
    try:
        notes_slide = slide.notes_slide
        notes_text = notes_slide.notes_text_frame.text
    except:
        pass
    check(len(notes_text) > 50, f"Has speaker notes ({len(notes_text)} chars)")
    check("CFO" in notes_text, "Speaker notes mention CFO")
    
    # Check Q4 dip addressed in notes
    if slide_num in [1, 2, 3, 6]:
        dip_keywords = ["dip", "decline", "Q4", "deferred", "transition", "shortfall"]
        has_dip_ref = any(kw.lower() in notes_text.lower() for kw in dip_keywords)
        check(has_dip_ref, "Notes address Q4 revenue context")

# --- 3. Bar chart on Slide 3 ---
print("\n=== BAR CHART (Slide 3) ===")
slide3 = prs.slides[2]
bar_charts = []
for shape in slide3.shapes:
    if shape.has_chart:
        ct = shape.chart.chart_type
        if ct in (XL_CHART_TYPE.COLUMN_CLUSTERED, XL_CHART_TYPE.BAR_CLUSTERED,
                  XL_CHART_TYPE.COLUMN_STACKED, XL_CHART_TYPE.BAR_STACKED):
            bar_charts.append(shape.chart)

check(len(bar_charts) >= 1, f"Slide 3 has a bar/column chart (found {len(bar_charts)})")

if bar_charts:
    chart = bar_charts[0]
    # Check data
    series = chart.series[0]
    values = [series.values[i] for i in range(len(series.values))]
    check(values == [2.4, 2.8, 3.1, 2.9], f"Chart data matches: {values}")
    
    # Check categories
    cats = [str(c) for c in chart.plots[0].categories]
    check(cats == ["Q1", "Q2", "Q3", "Q4"], f"Categories: {cats}")

# --- 4. Pie chart on Slide 5 ---
print("\n=== PIE CHART (Slide 5) ===")
slide5 = prs.slides[4]
pie_charts = []
for shape in slide5.shapes:
    if shape.has_chart:
        ct = shape.chart.chart_type
        if ct in (XL_CHART_TYPE.PIE, XL_CHART_TYPE.PIE_EXPLODED):
            pie_charts.append(shape.chart)

check(len(pie_charts) >= 1, f"Slide 5 has a pie chart (found {len(pie_charts)})")

if pie_charts:
    chart = pie_charts[0]
    series = chart.series[0]
    values = [series.values[i] for i in range(len(series.values))]
    check(sum(values) == 100, f"Pie values sum to 100 (got {sum(values)})")
    check(len(values) == 5, f"Pie has 5 categories (got {len(values)})")

# --- 5. Revenue data accuracy ---
print("\n=== DATA ACCURACY ===")
slide2_text = ""
for shape in prs.slides[1].shapes:
    if shape.has_text_frame:
        slide2_text += shape.text_frame.text
check("$11.2M" in slide2_text or "11.2" in slide2_text, "Full-year revenue $11.2M referenced")
check("4%" in slide2_text or "4 %" in slide2_text, "4% churn increase referenced")

# --- Summary ---
print(f"\n{'='*50}")
print(f"VERIFICATION SUMMARY: {PASS} passed, {FAIL} failed out of {PASS+FAIL} checks")
if FAIL == 0:
    print("🎉 ALL CHECKS PASSED!")
else:
    print(f"⚠️  {FAIL} check(s) failed — review above")
print(f"{'='*50}")
sys.exit(0 if FAIL == 0 else 1)

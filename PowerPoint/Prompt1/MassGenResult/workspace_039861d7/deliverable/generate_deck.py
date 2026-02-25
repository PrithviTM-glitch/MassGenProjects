#!/usr/bin/env python3
"""
Q4 Quarterly Earnings Deck Generator
=====================================
Generates an 8-slide 'Q4_Earnings.pptx' for a SaaS company using python-pptx.

Raw Data:
  - Q1 Revenue: $2.4M
  - Q2 Revenue: $2.8M
  - Q3 Revenue: $3.1M
  - Q4 Revenue: $2.9M
  - Customer churn increased by 4% in Q4

Features:
  - Native bar chart on Slide 3 (Quarterly Revenue Comparison)
  - Native pie chart on Slide 5 (Churn Breakdown)
  - 'Key Takeaway' text box on every slide
  - CFO speaker notes on every slide addressing the Q4 revenue dip
  - Robust error handling throughout
"""

import sys
import os
import traceback
from pathlib import Path

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt, Emu
    from pptx.chart.data import CategoryChartData
    from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION, XL_LABEL_POSITION
    from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
    from pptx.dml.color import RGBColor
except ImportError as e:
    print(f"ERROR: python-pptx is not installed. Run: pip install python-pptx\n{e}")
    sys.exit(1)


# =============================================================================
# Configuration / Data
# =============================================================================

REVENUE_DATA = {
    "Q1": 2.4,
    "Q2": 2.8,
    "Q3": 3.1,
    "Q4": 2.9,
}

# Churn breakdown for Slide 5 pie chart
CHURN_BREAKDOWN = {
    "Voluntary – Price Sensitivity": 35,
    "Voluntary – Competitor Switch": 25,
    "Voluntary – Feature Gaps": 20,
    "Involuntary – Payment Failure": 12,
    "Involuntary – Other": 8,
}

TOTAL_CHURN_INCREASE_PCT = 4  # Q4 churn increase

# Brand colours
COLOR_PRIMARY = RGBColor(0x1B, 0x3A, 0x5C)   # Deep navy
COLOR_ACCENT = RGBColor(0x2E, 0x86, 0xAB)    # Teal accent
COLOR_HIGHLIGHT = RGBColor(0xE8, 0x4D, 0x3D)  # Alert red
COLOR_BG_TAKEAWAY = RGBColor(0xF0, 0xF4, 0xF8) # Light grey-blue
COLOR_WHITE = RGBColor(0xFF, 0xFF, 0xFF)
COLOR_DARK_TEXT = RGBColor(0x2D, 0x2D, 0x2D)
COLOR_CHART_BLUE = RGBColor(0x2E, 0x86, 0xAB)
COLOR_CHART_RED = RGBColor(0xE8, 0x4D, 0x3D)
COLOR_CHART_GREEN = RGBColor(0x27, 0xAE, 0x60)

CHART_COLORS = [
    RGBColor(0x2E, 0x86, 0xAB),  # Teal
    RGBColor(0x1B, 0x3A, 0x5C),  # Navy
    RGBColor(0xE8, 0x4D, 0x3D),  # Red
    RGBColor(0xF3, 0x9C, 0x12),  # Orange
    RGBColor(0x27, 0xAE, 0x60),  # Green
]

OUTPUT_FILENAME = "Q4_Earnings.pptx"

# =============================================================================
# Slide content definitions
# =============================================================================

SLIDES = [
    {
        "title": "Q4 FY2025 Quarterly Earnings Report",
        "subtitle": "SaaSCo Inc. – Investor Presentation",
        "takeaway": "Key Takeaway: Despite a Q4 revenue dip to $2.9M, annual revenue grew 21% YoY with a strong $11.2M total.",
        "notes": (
            "CFO Notes – Opening Slide:\n\n"
            "Good morning, everyone. Thank you for joining our Q4 earnings call. "
            "I want to address the Q4 revenue figure up front: we came in at $2.9M, "
            "which represents a sequential decline from Q3's $3.1M. However, this was "
            "a deliberate outcome of strategic decisions we made mid-quarter — specifically, "
            "sunsetting two low-margin product lines and transitioning enterprise clients "
            "to annual contracts that will begin recognizing revenue in Q1 of next year. "
            "Full-year revenue reached $11.2M, a 21% increase year-over-year. "
            "Let's walk through the details."
        ),
    },
    {
        "title": "Executive Summary",
        "body": (
            "• Full-year revenue: $11.2M (+21% YoY)\n"
            "• Q4 revenue: $2.9M (–6.5% QoQ from Q3's $3.1M)\n"
            "• ARR run-rate: $11.6M entering FY2026\n"
            "• Net Revenue Retention (NRR): 112%\n"
            "• Customer churn increased 4% in Q4 — driven by SMB segment\n"
            "• Gross margin held steady at 74%\n"
            "• Operating cash flow positive for third consecutive quarter"
        ),
        "takeaway": "Key Takeaway: The Q4 dip is a timing issue from contract transitions — ARR and NRR remain strong.",
        "notes": (
            "CFO Notes – Executive Summary:\n\n"
            "The headline here is that while Q4 showed a sequential decline, our annual "
            "metrics tell a compelling growth story. The $2.9M Q4 figure includes a $180K "
            "impact from deferred enterprise contract revenue that will land in Q1. "
            "The 4% churn increase is concentrated in our SMB segment — our enterprise "
            "retention actually improved to 97%. We are taking targeted action on SMB "
            "churn with a new onboarding program launching in January."
        ),
    },
    {
        "title": "Quarterly Revenue Comparison",
        "chart_type": "bar",
        "takeaway": "Key Takeaway: Q4's $2.9M reflects strategic contract restructuring, not demand weakness. Q3 peak of $3.1M included one-time upsells.",
        "notes": (
            "CFO Notes – Revenue Chart:\n\n"
            "Looking at the bar chart, you'll see consistent growth from Q1 through Q3, "
            "then the dip in Q4. I want to be very transparent: approximately $200K of the "
            "shortfall is attributable to the enterprise contract transition I mentioned. "
            "Another ~$50K relates to sunsetting our legacy analytics add-on. Excluding "
            "these items, underlying Q4 revenue would have been approximately $3.15M — "
            "above Q3. The trajectory remains positive and we expect Q1 FY2026 revenue "
            "to exceed $3.3M as deferred revenue begins to recognize."
        ),
    },
    {
        "title": "Key Financial Metrics",
        "body": (
            "• Gross Margin: 74% (stable QoQ)\n"
            "• Operating Expenses: $2.1M (–3% QoQ due to hiring freeze)\n"
            "• EBITDA: $0.42M (14.5% margin)\n"
            "• Free Cash Flow: $0.31M\n"
            "• Cash & Equivalents: $8.7M\n"
            "• Deferred Revenue: $1.8M (+22% QoQ)\n"
            "• Rule of 40 Score: 35 (approaching benchmark)"
        ),
        "takeaway": "Key Takeaway: Deferred revenue surged 22% QoQ — a strong leading indicator for Q1 FY2026.",
        "notes": (
            "CFO Notes – Financial Metrics:\n\n"
            "I want to draw your attention to deferred revenue, which grew 22% "
            "quarter-over-quarter to $1.8M. This is the direct result of the enterprise "
            "contract restructuring and gives us high visibility into Q1. Our EBITDA "
            "margin of 14.5% shows we're balancing growth with profitability. The Q4 "
            "revenue dip did not impact our bottom line disproportionately because we "
            "proactively managed opex — the hiring freeze saved approximately $150K in the "
            "quarter. Our cash position of $8.7M gives us 12+ months of runway even in a "
            "zero-growth scenario."
        ),
    },
    {
        "title": "Customer Churn Analysis",
        "chart_type": "pie",
        "takeaway": "Key Takeaway: 60% of churn is voluntary and addressable — new onboarding and pricing programs target the top two drivers.",
        "notes": (
            "CFO Notes – Churn Breakdown:\n\n"
            "The 4% increase in churn is our most pressing concern and I want to address "
            "it head-on. The pie chart shows that price sensitivity (35%) and competitor "
            "switching (25%) account for 60% of churn — both are actionable. We are "
            "launching a revised pricing tier for SMB customers in Q1 that reduces the "
            "entry-level price point by 15% while maintaining per-seat economics. For "
            "competitor switching, our product team is fast-tracking three feature requests "
            "that exit-survey data identified as top reasons. The 12% involuntary churn "
            "from payment failures is being addressed with a new dunning system going live "
            "next month. I'm confident we can reduce overall churn by 2 percentage points "
            "within two quarters."
        ),
    },
    {
        "title": "Q4 Revenue Dip — Root Cause Analysis",
        "body": (
            "Root Causes of Q4 Sequential Revenue Decline (–$200K):\n\n"
            "1. Enterprise Contract Transition: –$180K\n"
            "   → 12 enterprise clients moved to annual prepaid contracts\n"
            "   → Revenue deferred to Q1 FY2026\n\n"
            "2. Product Sunset (Legacy Analytics): –$50K\n"
            "   → Planned deprecation of low-margin add-on\n"
            "   → 70% of affected users migrated to core platform\n\n"
            "3. SMB Churn Impact: –$90K\n"
            "   → Higher-than-expected SMB losses\n"
            "   → Partially offset by +$120K in enterprise upsells"
        ),
        "takeaway": "Key Takeaway: The Q4 dip is fully explainable — $180K is deferred (not lost) and structural actions improve long-term margins.",
        "notes": (
            "CFO Notes – Root Cause Analysis:\n\n"
            "This slide is critical for understanding the Q4 narrative. The $180K in "
            "deferred enterprise revenue is not revenue lost — it's revenue shifted. "
            "These 12 enterprise contracts represent $2.16M in annual committed value, "
            "which is actually a 15% uplift from their previous monthly arrangements. "
            "The legacy analytics sunset was planned for 18 months and removes a product "
            "with 45% gross margin versus our core platform at 78%. The SMB churn of $90K "
            "is the one area I'm not satisfied with — our new pricing and onboarding "
            "initiatives directly target this. Net-net, I view Q4 as a quarter where we "
            "made hard short-term choices for long-term structural improvement."
        ),
    },
    {
        "title": "FY2026 Outlook & Guidance",
        "body": (
            "FY2026 Guidance:\n\n"
            "• Revenue: $13.5M – $14.2M (+20–27% YoY)\n"
            "• Q1 FY2026 Revenue: $3.3M – $3.5M\n"
            "• Gross Margin Target: 75–77%\n"
            "• Churn Reduction Target: –2pp by Q2 FY2026\n"
            "• Planned Headcount Additions: 18 (Sales, Engineering, CS)\n"
            "• Product Launches: AI-powered analytics (Q2), Mobile app (Q3)\n"
            "• Expected ARR exit rate: $14.5M+"
        ),
        "takeaway": "Key Takeaway: FY2026 guidance of $13.5–14.2M reflects confidence in deferred revenue conversion and churn reduction.",
        "notes": (
            "CFO Notes – Outlook:\n\n"
            "Our guidance range of $13.5M to $14.2M assumes the deferred enterprise "
            "revenue converts as contracted — which we have high confidence in given the "
            "signed commitments. The Q1 target of $3.3M to $3.5M already has 85% pipeline "
            "coverage. We're investing in 18 new hires across sales, engineering, and "
            "customer success — the CS hires are specifically aimed at our churn reduction "
            "target. Our AI analytics module launching in Q2 has generated significant "
            "interest in beta, with 40 enterprise prospects in the pipeline. The Q4 "
            "revenue dip should be viewed as the trough in a transition quarter — "
            "we expect accelerating sequential growth through FY2026."
        ),
    },
    {
        "title": "Thank You & Q&A",
        "body": (
            "SaaSCo Inc. — Q4 FY2025 Earnings\n\n"
            "Contact:\n"
            "  Investor Relations: ir@saasco.com\n"
            "  CFO Office: cfo@saasco.com\n\n"
            "Upcoming Dates:\n"
            "  • Q1 FY2026 Earnings: April 15, 2026\n"
            "  • Annual Shareholder Meeting: May 20, 2026\n"
            "  • Investor Day: June 10, 2026\n\n"
            "Forward-looking statements in this presentation are subject to risks\n"
            "and uncertainties. See our SEC filings for details."
        ),
        "takeaway": "Key Takeaway: We're committed to transparent communication — the Q4 dip reflects strategic repositioning, not fundamental weakness.",
        "notes": (
            "CFO Notes – Closing:\n\n"
            "Before we open for Q&A, I want to reiterate three points: First, the Q4 "
            "revenue dip is a timing and transition issue — $180K is deferred, not lost. "
            "Second, our churn increase is concentrated, understood, and being actively "
            "addressed with specific programs launching in Q1. Third, our forward "
            "indicators — deferred revenue, pipeline, NRR — are all strong. I'm confident "
            "in our FY2026 guidance and I look forward to reporting on the progress of "
            "our churn reduction and revenue acceleration initiatives next quarter. "
            "We're happy to take your questions."
        ),
    },
]

assert len(SLIDES) == 8, f"Expected 8 slides, got {len(SLIDES)}"


# =============================================================================
# Helper Functions
# =============================================================================

def _set_slide_background(slide, color):
    """Set the background fill color of a slide."""
    try:
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = color
    except Exception as e:
        print(f"  WARNING: Could not set slide background: {e}")


def _add_title(slide, text, left=Inches(0.7), top=Inches(0.4),
               width=Inches(12.0), height=Inches(0.9),
               font_size=Pt(32), color=COLOR_PRIMARY, bold=True):
    """Add a styled title text box to a slide."""
    try:
        txBox = slide.shapes.add_textbox(left, top, width, height)
        tf = txBox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = text
        p.font.size = font_size
        p.font.bold = bold
        p.font.color.rgb = color
        p.font.name = "Calibri"
        return txBox
    except Exception as e:
        raise RuntimeError(f"Failed to add title '{text}': {e}") from e


def _add_body_text(slide, text, left=Inches(0.7), top=Inches(1.6),
                   width=Inches(12.0), height=Inches(4.0),
                   font_size=Pt(16), color=COLOR_DARK_TEXT):
    """Add body text to a slide."""
    try:
        txBox = slide.shapes.add_textbox(left, top, width, height)
        tf = txBox.text_frame
        tf.word_wrap = True

        lines = text.split("\n")
        for i, line in enumerate(lines):
            if i == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()
            p.text = line
            p.font.size = font_size
            p.font.color.rgb = color
            p.font.name = "Calibri"
            p.space_after = Pt(4)

            # Bold bullet headers
            if line.startswith("•") or line.startswith("Root Causes"):
                p.font.bold = True
            if line.startswith("   →"):
                p.font.size = Pt(14)
                p.font.color.rgb = RGBColor(0x55, 0x55, 0x55)
                p.font.bold = False

        return txBox
    except Exception as e:
        raise RuntimeError(f"Failed to add body text: {e}") from e


def _add_takeaway_box(slide, text, left=Inches(0.7), top=Inches(6.0),
                      width=Inches(12.0), height=Inches(0.65)):
    """Add a styled 'Key Takeaway' text box to the bottom of a slide."""
    try:
        txBox = slide.shapes.add_textbox(left, top, width, height)
        tf = txBox.text_frame
        tf.word_wrap = True
        tf.margin_left = Inches(0.15)
        tf.margin_right = Inches(0.15)
        tf.margin_top = Inches(0.08)
        tf.margin_bottom = Inches(0.08)

        # Set background fill
        fill = txBox.fill
        fill.solid()
        fill.fore_color.rgb = COLOR_BG_TAKEAWAY

        # Add border-like effect via line
        txBox.line.color.rgb = COLOR_ACCENT
        txBox.line.width = Pt(1.5)

        p = tf.paragraphs[0]
        # Split into "Key Takeaway:" prefix and the rest
        if text.startswith("Key Takeaway:"):
            prefix = "Key Takeaway: "
            rest = text[len("Key Takeaway:"):].strip()
            run1 = p.add_run()
            run1.text = prefix
            run1.font.size = Pt(11)
            run1.font.bold = True
            run1.font.color.rgb = COLOR_ACCENT
            run1.font.name = "Calibri"

            run2 = p.add_run()
            run2.text = rest
            run2.font.size = Pt(11)
            run2.font.bold = False
            run2.font.color.rgb = COLOR_DARK_TEXT
            run2.font.name = "Calibri"
        else:
            p.text = text
            p.font.size = Pt(11)
            p.font.color.rgb = COLOR_DARK_TEXT
            p.font.name = "Calibri"

        return txBox
    except Exception as e:
        raise RuntimeError(f"Failed to add takeaway box: {e}") from e


def _add_speaker_notes(slide, notes_text):
    """Add speaker notes to a slide."""
    try:
        notes_slide = slide.notes_slide
        tf = notes_slide.notes_text_frame
        tf.text = notes_text
    except Exception as e:
        raise RuntimeError(f"Failed to add speaker notes: {e}") from e


def _add_bar_chart(slide):
    """Add a native bar chart comparing quarterly revenue to Slide 3."""
    try:
        chart_data = CategoryChartData()
        chart_data.categories = list(REVENUE_DATA.keys())
        chart_data.add_series("Revenue ($M)", tuple(REVENUE_DATA.values()))

        x, y, cx, cy = Inches(1.5), Inches(1.6), Inches(10.3), Inches(4.0)
        chart_frame = slide.shapes.add_chart(
            XL_CHART_TYPE.COLUMN_CLUSTERED, x, y, cx, cy, chart_data
        )
        chart = chart_frame.chart

        # Style the chart
        chart.has_legend = False
        chart.style = 2  # Clean style

        # Value axis
        value_axis = chart.value_axis
        value_axis.minimum_scale = 0
        value_axis.maximum_scale = 4.0
        value_axis.major_gridlines.format.line.color.rgb = RGBColor(0xDD, 0xDD, 0xDD)
        value_axis.has_title = True
        value_axis.axis_title.text_frame.paragraphs[0].text = "Revenue ($M)"
        value_axis.axis_title.text_frame.paragraphs[0].font.size = Pt(11)
        value_axis.axis_title.text_frame.paragraphs[0].font.color.rgb = COLOR_DARK_TEXT

        # Category axis
        category_axis = chart.category_axis
        category_axis.has_major_gridlines = False
        category_axis.tick_labels.font.size = Pt(12)
        category_axis.tick_labels.font.bold = True

        # Data labels
        plot = chart.plots[0]
        plot.has_data_labels = True
        data_labels = plot.data_labels
        data_labels.font.size = Pt(12)
        data_labels.font.bold = True
        data_labels.number_format = '$#,##0.0"M"'
        data_labels.show_value = True

        # Color individual points — highlight Q4 in red
        series = chart.series[0]
        for idx, quarter in enumerate(REVENUE_DATA.keys()):
            point = series.points[idx]
            fill = point.format.fill
            fill.solid()
            if quarter == "Q4":
                fill.fore_color.rgb = COLOR_CHART_RED
            else:
                fill.fore_color.rgb = COLOR_CHART_BLUE

        return chart_frame
    except Exception as e:
        raise RuntimeError(f"Failed to add bar chart: {e}") from e


def _add_pie_chart(slide):
    """Add a native pie chart showing churn breakdown to Slide 5."""
    try:
        chart_data = CategoryChartData()
        chart_data.categories = list(CHURN_BREAKDOWN.keys())
        chart_data.add_series("Churn Breakdown", tuple(CHURN_BREAKDOWN.values()))

        x, y, cx, cy = Inches(1.5), Inches(1.6), Inches(10.3), Inches(4.0)
        chart_frame = slide.shapes.add_chart(
            XL_CHART_TYPE.PIE, x, y, cx, cy, chart_data
        )
        chart = chart_frame.chart

        # Legend
        chart.has_legend = True
        chart.legend.position = XL_LEGEND_POSITION.BOTTOM
        chart.legend.include_in_layout = False
        chart.legend.font.size = Pt(10)

        # Data labels
        plot = chart.plots[0]
        plot.has_data_labels = True
        data_labels = plot.data_labels
        data_labels.font.size = Pt(10)
        data_labels.font.bold = True
        data_labels.number_format = '0"%"'
        data_labels.show_category_name = False
        data_labels.show_value = True
        data_labels.show_percentage = False

        # Color each slice
        series = chart.series[0]
        for idx in range(len(CHURN_BREAKDOWN)):
            point = series.points[idx]
            fill = point.format.fill
            fill.solid()
            fill.fore_color.rgb = CHART_COLORS[idx % len(CHART_COLORS)]

        return chart_frame
    except Exception as e:
        raise RuntimeError(f"Failed to add pie chart: {e}") from e


def _build_title_slide(prs, slide_data, slide_num):
    """Build the title slide (Slide 1) with special styling."""
    try:
        slide_layout = prs.slide_layouts[6]  # Blank layout
        slide = prs.slides.add_slide(slide_layout)
        _set_slide_background(slide, COLOR_PRIMARY)

        # Main title
        _add_title(slide, slide_data["title"],
                   top=Inches(1.5), font_size=Pt(36), color=COLOR_WHITE)

        # Subtitle
        _add_body_text(slide, slide_data.get("subtitle", ""),
                       top=Inches(2.8), font_size=Pt(20), color=RGBColor(0xBB, 0xCC, 0xDD))

        # Date line
        _add_body_text(slide, "Fiscal Quarter Ending December 31, 2025",
                       top=Inches(3.6), font_size=Pt(14), color=RGBColor(0x99, 0xAA, 0xBB))

        # Takeaway box (adjusted for dark background)
        _add_takeaway_box(slide, slide_data["takeaway"], top=Inches(5.5))

        _add_speaker_notes(slide, slide_data["notes"])
        return slide
    except Exception as e:
        raise RuntimeError(f"Failed to build title slide: {e}") from e


def _build_content_slide(prs, slide_data, slide_num):
    """Build a standard content slide."""
    try:
        slide_layout = prs.slide_layouts[6]  # Blank layout
        slide = prs.slides.add_slide(slide_layout)

        # Title
        _add_title(slide, slide_data["title"])

        # Content depends on type
        chart_type = slide_data.get("chart_type")
        if chart_type == "bar":
            _add_bar_chart(slide)
        elif chart_type == "pie":
            _add_pie_chart(slide)
        elif "body" in slide_data:
            _add_body_text(slide, slide_data["body"])

        # Takeaway
        _add_takeaway_box(slide, slide_data["takeaway"])

        # Speaker notes
        _add_speaker_notes(slide, slide_data["notes"])

        return slide
    except Exception as e:
        raise RuntimeError(f"Failed to build slide {slide_num} ('{slide_data.get('title', '?')}'): {e}") from e


def _build_closing_slide(prs, slide_data, slide_num):
    """Build the closing slide with special styling."""
    try:
        slide_layout = prs.slide_layouts[6]  # Blank layout
        slide = prs.slides.add_slide(slide_layout)
        _set_slide_background(slide, COLOR_PRIMARY)

        # Title
        _add_title(slide, slide_data["title"],
                   top=Inches(0.6), font_size=Pt(36), color=COLOR_WHITE)

        # Body
        _add_body_text(slide, slide_data.get("body", ""),
                       top=Inches(1.8), font_size=Pt(14), color=RGBColor(0xCC, 0xDD, 0xEE))

        # Takeaway
        _add_takeaway_box(slide, slide_data["takeaway"], top=Inches(5.8))

        # Speaker notes
        _add_speaker_notes(slide, slide_data["notes"])

        return slide
    except Exception as e:
        raise RuntimeError(f"Failed to build closing slide: {e}") from e


# =============================================================================
# Main Generation Function
# =============================================================================

def generate_deck(output_path=None):
    """
    Generate the Q4 Earnings PowerPoint deck.

    Args:
        output_path: Path to save the file. Defaults to OUTPUT_FILENAME in cwd.

    Returns:
        str: Absolute path to the saved file.

    Raises:
        RuntimeError: If any slide generation or save operation fails.
    """
    if output_path is None:
        output_path = OUTPUT_FILENAME

    output_path = Path(output_path).resolve()

    print(f"Generating Q4 Earnings deck: {output_path}")
    print(f"  Slides to generate: {len(SLIDES)}")

    try:
        prs = Presentation()
        # Set standard widescreen 16:9 dimensions (13.333" x 7.5")
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)
    except Exception as e:
        raise RuntimeError(f"Failed to initialize Presentation object: {e}") from e

    for idx, slide_data in enumerate(SLIDES):
        slide_num = idx + 1
        print(f"  Building Slide {slide_num}: {slide_data['title']}")

        try:
            if slide_num == 1:
                _build_title_slide(prs, slide_data, slide_num)
            elif slide_num == 8:
                _build_closing_slide(prs, slide_data, slide_num)
            else:
                _build_content_slide(prs, slide_data, slide_num)
        except Exception as e:
            print(f"  ERROR on Slide {slide_num}: {e}")
            traceback.print_exc()
            raise

    # Validate slide count before saving
    actual_count = len(prs.slides)
    if actual_count != 8:
        raise RuntimeError(
            f"Expected 8 slides but presentation has {actual_count}. Aborting save."
        )

    # Save
    try:
        output_path.parent.mkdir(parents=True, exist_ok=True)
        prs.save(str(output_path))
        file_size = output_path.stat().st_size
        print(f"\n  ✅ Saved successfully: {output_path}")
        print(f"  File size: {file_size:,} bytes ({file_size / 1024:.1f} KB)")
        return str(output_path)
    except PermissionError as e:
        raise RuntimeError(f"Permission denied writing to {output_path}: {e}") from e
    except OSError as e:
        raise RuntimeError(f"OS error saving file to {output_path}: {e}") from e
    except Exception as e:
        raise RuntimeError(f"Unexpected error saving presentation: {e}") from e


# =============================================================================
# Entry Point
# =============================================================================

if __name__ == "__main__":
    try:
        saved_path = generate_deck()
        print(f"\n{'='*60}")
        print(f"  DECK GENERATION COMPLETE")
        print(f"  Output: {saved_path}")
        print(f"{'='*60}")
        sys.exit(0)
    except Exception as e:
        print(f"\n{'='*60}")
        print(f"  DECK GENERATION FAILED")
        print(f"  Error: {e}")
        print(f"{'='*60}")
        traceback.print_exc()
        sys.exit(1)

#!/usr/bin/env python3
"""Comprehensive verification of Q4_Earnings.pptx"""

from pptx import Presentation
from pptx.enum.chart import XL_CHART_TYPE
import sys

PPTX_PATH = "deliverable/Q4_Earnings.pptx"

passed = 0
failed = 0

def check(condition, desc):
    global passed, failed
    if condition:
        passed += 1
        print(f"  ✅ {desc}")
    else:
        failed += 1
        print(f"  ❌ {desc}")

prs = Presentation(PPTX_PATH)

# 1. Slide count
check(len(prs.slides) == 8, f"Slide count is 8 (got {len(prs.slides)})")

# 2. Slide dimensions (16:9)
w_in = prs.slide_width / 914400
h_in = prs.slide_height / 914400
check(abs(w_in - 13.333) < 0.01, f"Slide width ~13.333\" for 16:9 (got {w_in:.3f}\")")
check(abs(h_in - 7.5) < 0.01, f"Slide height 7.5\" (got {h_in:.3f}\")")

# Per-slide checks
expected_titles = [
    "Q4 FY2025 Quarterly Earnings Report",
    "Executive Summary",
    "Quarterly Revenue Comparison",
    "Key Financial Metrics",
    "Customer Churn Analysis",
    "Q4 Revenue Dip — Root Cause Analysis",
    "FY2026 Outlook & Guidance",
    "Thank You & Q&A",
]

for i, slide in enumerate(prs.slides):
    slide_num = i + 1
    print(f"\n--- Slide {slide_num} ---")

    # Find all text content
    all_text = ""
    for shape in slide.shapes:
        if shape.has_text_frame:
            for para in shape.text_frame.paragraphs:
                all_text += para.text + "\n"

    # Title check
    check(expected_titles[i] in all_text, f"Title found: '{expected_titles[i][:50]}'")

    # Key Takeaway check
    check("Key Takeaway" in all_text, "Has 'Key Takeaway' text box")

    # Speaker notes check
    notes_text = ""
    try:
        notes_text = slide.notes_slide.notes_text_frame.text
    except:
        pass
    check(len(notes_text) > 50, f"Has speaker notes ({len(notes_text)} chars)")
    check("CFO" in notes_text, "Notes contain 'CFO' reference")

    # Chart checks for specific slides
    charts_on_slide = []
    for shape in slide.shapes:
        if shape.has_chart:
            charts_on_slide.append(shape.chart.chart_type)

    if slide_num == 3:
        check(len(charts_on_slide) > 0, "Slide 3 has a chart")
        if charts_on_slide:
            check(charts_on_slide[0] == XL_CHART_TYPE.COLUMN_CLUSTERED,
                  f"Slide 3 chart is COLUMN_CLUSTERED (bar chart)")
            # Verify data values
            chart = None
            for shape in slide.shapes:
                if shape.has_chart:
                    chart = shape.chart
                    break
            if chart:
                series = chart.series[0]
                values = [series.values[j] for j in range(len(series.values))]
                check(values == [2.4, 2.8, 3.1, 2.9],
                      f"Revenue data correct: {values}")

    if slide_num == 5:
        check(len(charts_on_slide) > 0, "Slide 5 has a chart")
        if charts_on_slide:
            check(charts_on_slide[0] == XL_CHART_TYPE.PIE,
                  f"Slide 5 chart is PIE (pie chart)")

# Q4 dip in notes
print("\n--- CFO Notes Q4 Dip Coverage ---")
dip_mentions = 0
for i, slide in enumerate(prs.slides):
    try:
        notes = slide.notes_slide.notes_text_frame.text.lower()
        if any(term in notes for term in ["dip", "decline", "shortfall", "$2.9m", "sequential"]):
            dip_mentions += 1
    except:
        pass
check(dip_mentions >= 3, f"At least 3 slides mention Q4 dip in notes (found {dip_mentions})")

print(f"\n{'='*50}")
print(f"RESULTS: {passed} passed, {failed} failed, {passed+failed} total")
print(f"{'='*50}")

sys.exit(0 if failed == 0 else 1)
